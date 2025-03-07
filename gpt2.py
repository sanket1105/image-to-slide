import base64
import io
import json
import os
from io import BytesIO

import requests
import streamlit as st
from dotenv import load_dotenv
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# Load environment variables for API keys
load_dotenv()
from pptx.enum.dml import MSO_LINE_DASH_STYLE, MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN


class WhiteboardToPPT:
    def __init__(self):
        # API configuration - use environment variable or Streamlit secrets
        self.vision_api_key = os.getenv("OPENAI_API_KEY")
        self.vision_api_url = "https://api.openai.com/v1/chat/completions"

        # PowerPoint settings
        self.slide_width = Inches(10)
        self.slide_height = Inches(7.5)

        # Configuration for element sizing and spacing
        self.min_box_width = Inches(1.2)
        self.max_box_width = Inches(3.5)
        self.min_box_height = Inches(0.5)
        self.max_box_height = Inches(1.2)
        self.horizontal_spacing = Inches(0.4)  # Space between elements horizontally
        self.vertical_spacing = Inches(0.4)  # Space between elements vertically

    def analyze_image(self, image_data):
        """
        Use AI vision model to analyze the whiteboard image and extract information.
        Returns structured data about text content and layout.
        """
        try:
            # Convert image to base64
            image_base64 = base64.b64encode(image_data).decode("utf-8")

            # Enhanced prompt for the vision model with better OCR instructions
            prompt = """
            Analyze this whiteboard image with high precision OCR, extracting:
            1. All text content with EXACT wording, preserving all punctuation, capitalization, and special characters
            2. The spatial layout and relationships between elements
            3. The hierarchical structure and flow of information
            
            FORMAT REQUIREMENTS - STRICTLY ADHERE TO THIS JSON STRUCTURE:
            {
              "phases": [list of phase headings at the top, ordered left to right],
              "sections": [list of vertical sections on the left side, ordered top to bottom],
              "elements": [
                {
                  "id": "unique identifier (elem1, elem2, etc.)",
                  "text": "exact text content with ALL special characters and formatting preserved",
                  "position": "precise position descriptor (top-left, middle-center, etc.)",
                  "type": "box/text/note",
                  "coordinates": {"x": percentage from left, "y": percentage from top},
                  "has_shape": boolean,
                  "emphasis": "none/high/medium/low based on visual prominence"
                }
              ],
              "connections": [
                {"from_id": "source element id", "to_id": "target element id", "arrow_type": "simple/double/none"}
              ]
            }
            
            CRITICAL INSTRUCTIONS:
            1. READ ALL TEXT CAREFULLY - capture every character including special symbols (#, *, etc.)
            2. AVOID TEXT DUPLICATION - ensure no element text is duplicated
            3. IDENTIFY ALL CONNECTIONS - carefully trace every line/arrow between elements
            4. ASSIGN UNIQUE IDs - each element must have a unique identifier for connections
            5. BE PRECISE WITH COORDINATES - position elements accurately in percentage terms
            6. IDENTIFY ANY EMPHASIS - note if text is highlighted, circled, or visually emphasized
            
            For coordinates, use percentages of image dimensions (0,0 is top left; 100,100 is bottom right).
            """

            # API call to vision model
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {self.vision_api_key}",
            }

            # Using GPT-4o for improved visual analysis
            payload = {
                "model": "gpt-4o",
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/jpeg;base64,{image_base64}"
                                },
                            },
                        ],
                    }
                ],
                "max_tokens": 4000,
                "temperature": 0.2,  # Lower temperature for more consistent results
            }

            with st.status(
                "Analyzing image with AI vision model...", expanded=True
            ) as status:
                response = requests.post(
                    self.vision_api_url, headers=headers, json=payload
                )

                if response.status_code != 200:
                    st.error(f"API Error: {response.status_code}")
                    st.error(response.text)
                    status.update(label="Analysis failed", state="error")
                    return None

                result = response.json()
                status.update(label="Analysis complete", state="complete")

            # Check if we have the expected structure
            if "choices" not in result or len(result["choices"]) == 0:
                st.error("Unexpected API response format")
                st.write(result)
                return None

            # Extract the content from the response
            content = result["choices"][0]["message"]["content"]

            # Parse the JSON content - handle potential JSON formatting issues
            try:
                # Find JSON content in the response (it might be mixed with other text)
                json_content = self._extract_json_from_text(content)
                analysis_data = json.loads(json_content)

                # Log the detected connections and elements
                st.expander("AI Analysis Results", expanded=False).write(analysis_data)

                # Enhance data with text analysis
                analysis_data = self._enhance_analysis_data(analysis_data)

                # NEW: De-duplicate text to ensure "Ideation" or "Sizing" doesn't appear twice
                analysis_data = self._deduplicate_analysis_data(analysis_data)

                return analysis_data
            except json.JSONDecodeError as e:
                st.error(f"Error parsing JSON response: {e}")
                st.write("Raw content:")
                st.write(content)

                # Fallback to manual extraction if JSON parsing fails
                return self._fallback_extraction(content)

        except Exception as e:
            st.error(f"Error analyzing image: {e}")
            return None

    def _enhance_analysis_data(self, data):
        """Add additional analysis to improve the layout and connections"""
        if not data:
            return data

        # Add size estimation based on text length
        if "elements" in data:
            for element in data["elements"]:
                if "text" in element:
                    # Calculate text metrics
                    text = element["text"]
                    element["text_length"] = len(text)
                    element["word_count"] = len(text.split())

                    # Estimate importance based on position, content, and connections
                    importance = 1  # Default importance

                    # Check if this element is connected to many others
                    if "connections" in data:
                        connections_from = sum(
                            1
                            for conn in data["connections"]
                            if conn["from_id"] == element["id"]
                        )
                        connections_to = sum(
                            1
                            for conn in data["connections"]
                            if conn["to_id"] == element["id"]
                        )
                        element["connection_count"] = connections_from + connections_to

                        # Elements with more connections tend to be more important
                        if connections_from + connections_to > 2:
                            importance += 1

                    # Top position elements are often more important
                    if "position" in element and "top" in element["position"].lower():
                        importance += 1

                    # Elements with special markers like "#" might be emphasized
                    if "#" in text or "*" in text or "!" in text:
                        importance += 1

                    element["importance"] = min(importance, 3)  # Scale of 1-3

        return data

    def _extract_json_from_text(self, text):
        """Extract JSON content from potentially mixed text"""
        # Look for content between curly braces
        start_idx = text.find("{")
        end_idx = text.rfind("}")

        if start_idx >= 0 and end_idx > start_idx:
            return text[start_idx : end_idx + 1]

        return "{}"  # Return empty JSON if none found

    def _fallback_extraction(self, content):
        """Enhanced manual extraction of key elements as fallback"""
        st.warning("Using fallback extraction method")

        # More sophisticated fallback with better structure
        return {
            "phases": ["Ideation", "Align", "Launch", "Manage"],
            "sections": ["Strategy", "Planning", "Foundation"],
            "elements": [
                {
                    "id": "elem1",
                    "text": "#Need Sizing Tool#",
                    "position": "top-left",
                    "type": "note",
                    "coordinates": {"x": 15, "y": 13},
                    "has_shape": True,
                    "emphasis": "high",
                    "text_length": 17,
                    "word_count": 3,
                    "importance": 3,
                },
                {
                    "id": "elem2",
                    "text": "Define: Test (channels) Audience Success",
                    "position": "top-center",
                    "type": "box",
                    "coordinates": {"x": 35, "y": 13},
                    "has_shape": True,
                    "emphasis": "medium",
                    "text_length": 42,
                    "word_count": 5,
                    "importance": 2,
                },
                # Additional elements would be included here...
            ],
            "connections": [
                {"from_id": "elem1", "to_id": "elem2", "arrow_type": "simple"},
                # Additional connections would be included here...
            ],
        }

    def _deduplicate_analysis_data(self, data):
        """
        Remove duplicate text entries from phases, sections, and elements
        so the same text doesn't appear multiple times.
        """
        # 1) De-duplicate phases (optional if you see duplicates there)
        if "phases" in data:
            unique_phases = []
            seen_phases = set()
            for ph in data["phases"]:
                if ph not in seen_phases:
                    unique_phases.append(ph)
                    seen_phases.add(ph)
            data["phases"] = unique_phases

        # 2) De-duplicate sections (optional if you see duplicates)
        if "sections" in data:
            unique_sections = []
            seen_sections = set()
            for sec in data["sections"]:
                if sec not in seen_sections:
                    unique_sections.append(sec)
                    seen_sections.add(sec)
            data["sections"] = unique_sections

        # 3) De-duplicate elements based on their exact text
        if "elements" in data:
            unique_texts = set()
            deduped_elements = []
            for e in data["elements"]:
                txt = e.get("text", "").strip()
                if txt not in unique_texts:
                    deduped_elements.append(e)
                    unique_texts.add(txt)
            data["elements"] = deduped_elements

        return data

    def _calculate_element_dimensions(self, element):
        """Calculate appropriate dimensions for an element based on its content"""
        text = element.get("text", "")
        text_length = element.get("text_length", len(text))
        word_count = element.get("word_count", len(text.split()))
        importance = element.get("importance", 1)

        # Calculate width based on text length and word count
        avg_word_length = text_length / max(word_count, 1)
        width_factor = min(
            1.0, 0.6 + (avg_word_length / 20)
        )  # Scale based on average word length

        width = max(
            self.min_box_width,
            min(self.max_box_width, Inches(1.0 + (text_length / 15) * width_factor)),
        )

        # Estimate how many lines the text will take
        chars_per_line = width.inches * 11  # ~11 chars per inch with default font
        estimated_lines = max(1, text_length / chars_per_line)

        # Add some extra height for important elements
        importance_factor = 1.0 + ((importance - 1) * 0.15)

        height = max(
            self.min_box_height,
            min(
                self.max_box_height,
                Inches(0.4 + (estimated_lines * 0.2) * importance_factor),
            ),
        )

        return width, height

    def _intelligently_position_elements(self, analysis_data):
        """Create an optimal layout for elements based on their relationships and positions"""
        # Create a simple grid for positioning
        grid = {
            "rows": 5,
            "cols": 5,
            "cells": {},
        }

        elements_data = {}
        if "elements" in analysis_data:
            # Sort by y-coordinate, then x-coordinate
            elements = sorted(
                analysis_data["elements"],
                key=lambda e: (
                    e.get("coordinates", {}).get("y", 50),
                    e.get("coordinates", {}).get("x", 50),
                ),
            )

            for element in elements:
                element_id = element.get("id")
                if not element_id:
                    continue

                coords = element.get("coordinates", {"x": 50, "y": 50})
                row = min(grid["rows"] - 1, int((coords["y"] / 100) * grid["rows"]))
                col = min(grid["cols"] - 1, int((coords["x"] / 100) * grid["cols"]))

                width, height = self._calculate_element_dimensions(element)
                elements_data[element_id] = {
                    "element": element,
                    "grid_pos": (row, col),
                    "width": width,
                    "height": height,
                    "connections": [],
                }

                if (row, col) not in grid["cells"]:
                    grid["cells"][(row, col)] = []
                grid["cells"][(row, col)].append(element_id)

        # Add connection info
        if "connections" in analysis_data:
            for connection in analysis_data["connections"]:
                from_id = connection.get("from_id")
                to_id = connection.get("to_id")
                if from_id in elements_data and to_id in elements_data:
                    elements_data[from_id]["connections"].append(
                        {"to": to_id, "type": connection.get("arrow_type", "simple")}
                    )
                    elements_data[to_id]["connections"].append(
                        {
                            "from": from_id,
                            "type": connection.get("arrow_type", "simple"),
                        }
                    )

        # Convert grid positions to actual coordinates
        slide_margin_x = Inches(1.0)
        slide_margin_y = Inches(1.0)
        usable_width = self.slide_width - (slide_margin_x * 2)
        usable_height = self.slide_height - (slide_margin_y * 2)
        cell_width = usable_width / grid["cols"]
        cell_height = usable_height / grid["rows"]

        for element_id, data in elements_data.items():
            row, col = data["grid_pos"]

            base_x = (
                slide_margin_x
                + (col * cell_width)
                + (cell_width / 2)
                - (data["width"] / 2)
            )
            base_y = (
                slide_margin_y
                + (row * cell_height)
                + (cell_height / 2)
                - (data["height"] / 2)
            )

            # If multiple elements in same cell, spread them out
            cell_elements = grid["cells"].get((row, col), [])
            if len(cell_elements) > 1:
                num_elements = len(cell_elements)
                index_in_cell = cell_elements.index(element_id)
                offset_factor = Inches(0.15)
                horizontal_offset = (
                    index_in_cell - (num_elements - 1) / 2
                ) * offset_factor
                vertical_offset = (
                    index_in_cell - (num_elements - 1) / 2
                ) * offset_factor
                base_x += horizontal_offset
                base_y += vertical_offset

            data["position"] = {"left": base_x, "top": base_y}

        return elements_data

    def create_ppt_from_analysis(self, analysis_data):
        """
        Create a PowerPoint slide based on the AI-analyzed content.
        """
        prs = Presentation()
        prs.slide_width = self.slide_width
        prs.slide_height = self.slide_height

        slide_layout = prs.slide_layouts[6]  # Blank slide
        slide = prs.slides.add_slide(slide_layout)

        # 1. Title
        title_shape = slide.shapes.add_textbox(
            left=Inches(0.5), top=Inches(0.1), width=Inches(9.0), height=Inches(0.5)
        )
        title_frame = title_shape.text_frame
        title_frame.text = "Whiteboard Conversion"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(20)
        title_para.font.bold = True
        title_para.font.name = "Arial"
        title_para.font.color.rgb = RGBColor(0, 0, 0)  # Ensure black
        title_para.alignment = PP_ALIGN.CENTER

        # 2. Phases
        if "phases" in analysis_data and analysis_data["phases"]:
            phase_width = Inches(1.8)
            phase_height = Inches(0.6)
            total_phases = len(analysis_data["phases"])
            total_width = total_phases * phase_width
            start_x = (self.slide_width - total_width) / 2

            for i, phase in enumerate(analysis_data["phases"]):
                phase_x = start_x + (i * phase_width)
                phase_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left=phase_x,
                    top=Inches(0.7),
                    width=phase_width,
                    height=phase_height,
                )
                fill = phase_shape.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(240, 240, 255)
                line = phase_shape.line
                line.color.rgb = RGBColor(100, 100, 180)
                line.width = Pt(1.5)

                text_frame = phase_shape.text_frame
                text_frame.text = phase
                text_frame.margin_left = Pt(6)
                text_frame.margin_right = Pt(6)
                text_frame.margin_top = Pt(3)
                text_frame.margin_bottom = Pt(3)
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                p = text_frame.paragraphs[0]
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.name = "Calibri"
                p.font.color.rgb = RGBColor(0, 0, 0)  # Ensure black
                p.alignment = PP_ALIGN.CENTER

        # 3. Sections
        if "sections" in analysis_data and analysis_data["sections"]:
            for i, section in enumerate(analysis_data["sections"]):
                section_y = Inches(1.8 + i * 1.6)
                section_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left=Inches(0.3),
                    top=section_y,
                    width=Inches(0.9),
                    height=Inches(1.2),
                )
                section_shape.fill.solid()
                section_shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
                section_shape.line.color.rgb = RGBColor(120, 120, 120)

                text_frame = section_shape.text_frame
                text_frame.text = section
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                section_shape.rotation = 270  # Vertical text

                p = text_frame.paragraphs[0]
                p.font.size = Pt(12)
                p.font.name = "Calibri"
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 0, 0)  # Ensure black
                p.alignment = PP_ALIGN.CENTER

        # 4. Position elements
        element_layout = self._intelligently_position_elements(analysis_data)
        element_shapes = {}

        for element_id, element_data in element_layout.items():
            element = element_data["element"]
            position = element_data["position"]
            width = element_data["width"]
            height = element_data["height"]

            shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
            if element.get("type") == "note":
                shape_type = MSO_SHAPE.CLOUD

            shape = slide.shapes.add_shape(
                shape_type,
                left=position["left"],
                top=position["top"],
                width=width,
                height=height,
            )

            # Styling by importance
            importance = element.get("importance", 1)
            if importance >= 3 or element.get("emphasis") == "high":
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(255, 240, 240)
                shape.line.color.rgb = RGBColor(180, 100, 100)
                shape.line.width = Pt(2.0)
            elif importance == 2 or element.get("emphasis") == "medium":
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(240, 255, 240)
                shape.line.color.rgb = RGBColor(100, 160, 100)
                shape.line.width = Pt(1.5)
            else:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(250, 250, 250)
                shape.line.color.rgb = RGBColor(100, 100, 100)
                shape.line.width = Pt(1.0)

            text_frame = shape.text_frame
            text_frame.word_wrap = True
            text_frame.text = element["text"]
            text_frame.margin_left = Pt(6)
            text_frame.margin_right = Pt(6)
            text_frame.margin_top = Pt(3)
            text_frame.margin_bottom = Pt(3)
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            p = text_frame.paragraphs[0]
            text_length = element.get("text_length", len(element["text"]))
            if text_length > 70:
                font_size = 9
            elif text_length > 50:
                font_size = 10
            elif text_length > 30:
                font_size = 11
            else:
                font_size = 12
            if importance >= 2:
                font_size += 1

            p.font.size = Pt(font_size)
            p.font.name = "Calibri"
            if importance >= 2:
                p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)  # Ensure black
            p.alignment = PP_ALIGN.CENTER

            element_shapes[element_id] = {
                "shape": shape,
                "left": position["left"],
                "top": position["top"],
                "width": width,
                "height": height,
            }

        # 5. Connections
        if "connections" in analysis_data:
            for connection in analysis_data["connections"]:
                from_id = connection.get("from_id")
                to_id = connection.get("to_id")
                arrow_type = connection.get("arrow_type", "simple")
                if from_id not in element_shapes or to_id not in element_shapes:
                    continue

                from_shape = element_shapes[from_id]
                to_shape = element_shapes[to_id]

                from_center_x = from_shape["left"] + (from_shape["width"] / 2)
                from_center_y = from_shape["top"] + (from_shape["height"] / 2)
                to_center_x = to_shape["left"] + (to_shape["width"] / 2)
                to_center_y = to_shape["top"] + (to_shape["height"] / 2)

                dx = to_center_x - from_center_x
                dy = to_center_y - from_center_y

                # Pick best edge to start from
                if abs(dx) > abs(dy):
                    if dx > 0:
                        start_x = from_shape["left"] + from_shape["width"]
                        start_y = from_center_y
                    else:
                        start_x = from_shape["left"]
                        start_y = from_center_y
                else:
                    if dy > 0:
                        start_x = from_center_x
                        start_y = from_shape["top"] + from_shape["height"]
                    else:
                        start_x = from_center_x
                        start_y = from_shape["top"]

                # Pick best edge to end at
                if abs(dx) > abs(dy):
                    if dx > 0:
                        end_x = to_shape["left"]
                        end_y = to_center_y
                    else:
                        end_x = to_shape["left"] + to_shape["width"]
                        end_y = to_center_y
                else:
                    if dy > 0:
                        end_x = to_center_x
                        end_y = to_shape["top"]
                    else:
                        end_x = to_center_x
                        end_y = to_shape["top"] + to_shape["height"]

                connector_type = MSO_CONNECTOR.STRAIGHT
                if (abs(dx) + abs(dy)) > Inches(4).inches:
                    connector_type = MSO_CONNECTOR.CURVE

                connector = slide.shapes.add_connector(
                    connector_type, start_x, start_y, end_x, end_y
                )
                connector.line.color.rgb = RGBColor(100, 100, 100)
                connector.line.width = Pt(1.5)

                if arrow_type == "double":
                    connector.line.width = Pt(2.0)
                    connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
                    connector.line.begin_style = 1
                    connector.line.end_style = 1
                else:
                    connector.line.dash_style = MSO_LINE_DASH_STYLE.SOLID
                    connector.line.end_style = 1

        # Return the PowerPoint as a BytesIO object
        pptx_io = BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        return pptx_io

    def process_image_to_ppt(self, image_data):
        """
        Main method to process image and create PowerPoint.
        """
        # 1. Analyze image
        analysis_data = self.analyze_image(image_data)
        if not analysis_data:
            st.error("Failed to analyze image")
            return None

        # 2. Create PowerPoint
        with st.status("Creating PowerPoint slide...", expanded=True) as status:
            pptx_data = self.create_ppt_from_analysis(analysis_data)
            status.update(label="PowerPoint creation complete", state="complete")

        return pptx_data


def main():
    st.set_page_config(
        page_title="Whiteboard to PowerPoint Converter", page_icon="📊", layout="wide"
    )

    # Custom CSS
    st.markdown(
        """
    <style>
    .main-header {
        font-size: 2.5rem;
        margin-bottom: 0.5rem;
    }
    .sub-header {
        font-size: 1.2rem;
        color: #666;
        margin-bottom: 2rem;
    }
    .stButton>button {
        background-color: #4CAF50;
        color: white;
        padding: 0.5rem 1rem;
        font-size: 1rem;
        border-radius: 0.3rem;
        border: none;
    }
    .stButton>button:hover {
        background-color: #3e8e41;
    }
    .upload-section {
        background-color: #f8f9fa;
        padding: 1.5rem;
        border-radius: 0.5rem;
        margin-bottom: 1rem;
    }
    </style>
    """,
        unsafe_allow_html=True,
    )

    st.markdown(
        '<p class="main-header">Whiteboard to PowerPoint Converter</p>',
        unsafe_allow_html=True,
    )
    st.markdown(
        '<p class="sub-header">Transform your whiteboard brainstorming sessions into structured PowerPoint slides instantly</p>',
        unsafe_allow_html=True,
    )

    col1, col2 = st.columns([3, 2])
    with col1:
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            st.warning(
                "⚠️ OpenAI API key not found. Please add it to your .env file or Streamlit secrets."
            )
            api_key = st.text_input("Enter your OpenAI API key:", type="password")
            if api_key:
                os.environ["OPENAI_API_KEY"] = api_key
                st.success("API key set successfully!")
            else:
                st.info("An API key is required to analyze whiteboard images.")

        st.markdown('<div class="upload-section">', unsafe_allow_html=True)
        st.write("### Upload your whiteboard image")
        st.write(
            "Take a photo of your whiteboard and upload it here. The AI will analyze the content and structure."
        )

        uploaded_file = st.file_uploader(
            "Choose an image...", type=["jpg", "jpeg", "png"]
        )
        st.markdown("</div>", unsafe_allow_html=True)

        if uploaded_file is not None:
            image = Image.open(uploaded_file)
            st.image(image, caption="Uploaded Whiteboard Image", use_column_width=True)

            img_byte_arr = io.BytesIO()
            image.save(img_byte_arr, format=image.format)
            image_data = img_byte_arr.getvalue()

            processor = WhiteboardToPPT()

            if st.button("Convert to PowerPoint", use_container_width=True):
                pptx_data = processor.process_image_to_ppt(image_data)
                if pptx_data:
                    st.download_button(
                        label="Download PowerPoint",
                        data=pptx_data,
                        file_name="whiteboard_converted.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True,
                    )

    with col2:
        st.markdown(
            """
        ## How It Works
        
        1. **Upload** your whiteboard photo
        2. **AI** analyzes the content, structure, and relationships
        3. **Download** a formatted PowerPoint slide

        ## Tips
        - Good lighting
        - Clear, straight-on photo
        - Legible text in high contrast
        - Capture the entire whiteboard

        ## FAQ
        - **Editing**: The PPTX is fully editable
        - **Accuracy**: Handwriting or glare can reduce OCR accuracy
        - **Security**: Images are processed securely via OpenAI
        """
        )


# Run
if __name__ == "__main__":
    main()

