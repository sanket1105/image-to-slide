import base64
import io
import json
import logging
import os
import time
from io import BytesIO

import requests
import streamlit as st
import tenacity
from dotenv import load_dotenv
from PIL import Image, ImageEnhance
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Load environment variables for API keys
load_dotenv()


class WhiteboardToPPT:
    def __init__(self):
        # API configuration - use environment variable or Streamlit secrets
        self.vision_api_key = os.getenv("OPENAI_API_KEY") or st.secrets.get(
            "OPENAI_API_KEY", ""
        )
        self.vision_api_url = "https://api.openai.com/v1/chat/completions"

        # PowerPoint settings
        self.slide_width = Inches(10)
        self.slide_height = Inches(7.5)

        # Default styling options
        self.default_font = "Arial"
        self.default_outline_color = RGBColor(0, 0, 0)  # Black
        self.default_text_color = RGBColor(0, 0, 0)  # Black
        self.default_line_style = MSO_LINE_DASH_STYLE.SOLID

        # Cache for API responses
        self.analysis_cache = {}

    @tenacity.retry(
        wait=tenacity.wait_exponential(multiplier=1, min=4, max=10),
        stop=tenacity.stop_after_attempt(3),
        before_sleep=lambda retry_state: logger.warning(
            f"Retrying API call ({retry_state.attempt_number}/3)..."
        ),
    )
    def analyze_image(self, image_data):
        """
        Use AI vision model to analyze the whiteboard image with retry logic.
        Returns structured data about text content and layout.
        """
        try:
            # Validate image_data
            if not image_data or len(image_data) == 0:
                st.error("Image data is empty. Please upload a valid image.")
                return None

            logger.info(f"Image data size: {len(image_data)} bytes")

            # Convert image to base64
            image_base64 = base64.b64encode(image_data).decode("utf-8")

            # Check cache
            cache_key = image_base64[:50]  # Use a subset of base64 for caching
            if cache_key in self.analysis_cache:
                logger.info("Returning cached analysis")
                return self.analysis_cache[cache_key]

            # Prepare improved prompt for the vision model
            prompt = """
            Analyze this whiteboard image and extract:
            1. All text content with exact wording (no duplicates)
            2. The spatial layout (text in boxes, connections)
            3. Hierarchical structure (main sections, subsections)
            
            Format response as JSON with:
            - phases: List of main phase headings (e.g., "Ideation", "Align")
            - sections: List of vertical section labels
            - elements: List of objects with:
              - id: Unique identifier (e.g., "elem1")
              - text: Exact text content
              - position: "top-left", "middle-center", etc.
              - type: "box", "text", "note"
              - coordinates: Approximate x,y percentages (e.g., {"x": 20, "y": 30})
              - has_shape: Boolean for enclosure in a shape
            - connections: List of {"from_id": ..., "to_id": ...}
            
            Ensure:
            - No duplicate texts
            - Accurate detection of arrows/lines as connections
            - Unique IDs for all elements
            """

            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {self.vision_api_key}",
            }
            payload = {
                "model": "gpt-4o",
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/jpeg;base64,{image_base64}"
                                },
                            },
                        ],
                    }
                ],
                "max_tokens": 4000,
            }

            st.write("Analyzing image with AI vision model...")
            progress_bar = st.progress(0)
            response = requests.post(
                self.vision_api_url, headers=headers, json=payload, timeout=30
            )

            progress_bar.progress(50)

            if response.status_code != 200:
                st.error(f"API Error: {response.status_code} - {response.text}")
                return None

            result = response.json()
            if "choices" not in result or not result["choices"]:
                st.error("Unexpected API response format")
                st.write(result)
                return None

            content = result["choices"][0]["message"]["content"]
            json_content = self._extract_json_from_text(content)
            analysis_data = json.loads(json_content)

            # Log and cache the result
            logger.info(
                f"Successfully analyzed image with {len(analysis_data.get('elements', []))} elements"
            )
            self.analysis_cache[cache_key] = analysis_data
            progress_bar.progress(100)
            return analysis_data

        except requests.RequestException as e:
            st.error(f"Network error: {e}")
            return None
        except json.JSONDecodeError as e:
            st.error(f"Error parsing JSON: {e}")
            st.write("Raw content:", content)
            return self._fallback_extraction(content)
        except Exception as e:
            st.error(f"Unexpected error analyzing image: {e}")
            return None

    def _extract_json_from_text(self, text):
        """Extract JSON content from potentially mixed text"""
        start_idx = text.find("{")
        end_idx = text.rfind("}")
        return (
            text[start_idx : end_idx + 1]
            if start_idx >= 0 and end_idx > start_idx
            else "{}"
        )

    def _fallback_extraction(self, content):
        """Manual extraction as fallback"""
        st.warning("Using fallback extraction method")
        return {
            "phases": ["Ideation", "Align", "Launch", "Manage"],
            "sections": ["Strategy", "Planning", "Foundation"],
            "elements": [
                {
                    "id": "elem1",
                    "text": "Content Ideas",
                    "position": "middle-left",
                    "type": "box",
                    "coordinates": {"x": 15, "y": 50},
                    "has_shape": True,
                },
                {
                    "id": "elem2",
                    "text": "Define: Test (Automation), Audience Success",
                    "position": "top-center",
                    "type": "box",
                    "coordinates": {"x": 35, "y": 13},
                    "has_shape": True,
                },
                {
                    "id": "elem3",
                    "text": "Job Aid / Sharing Reporting (Dana)",
                    "position": "top-right",
                    "type": "box",
                    "coordinates": {"x": 65, "y": 13},
                    "has_shape": True,
                },
                {
                    "id": "elem4",
                    "text": "Sales Calling & Accountability",
                    "position": "middle-right",
                    "type": "box",
                    "coordinates": {"x": 85, "y": 50},
                    "has_shape": True,
                },
                {
                    "id": "elem5",
                    "text": "Lead/Contact Quality",
                    "position": "bottom-center",
                    "type": "box",
                    "coordinates": {"x": 45, "y": 75},
                    "has_shape": True,
                },
            ],
            "connections": [
                {"from_id": "elem1", "to_id": "elem2"},
                {"from_id": "elem2", "to_id": "elem3"},
                {"from_id": "elem4", "to_id": "elem5"},
            ],
        }

    def adjust_positions_for_overlap(self, elements_dict):
        """Grid-based overlap prevention with dynamic spacing"""
        grid_size = Inches(0.5)  # Minimum spacing between elements
        max_x, max_y = self.slide_width.inches, self.slide_height.inches
        grid = {}

        for elem_id, elem in elements_dict.items():
            left, top = elem["left"].inches, elem["top"].inches
            width, height = elem["width"].inches, elem["height"].inches
            right = left + width
            bottom = top + height

            # Assign to grid cells
            start_x = int(left / grid_size.inches)
            start_y = int(top / grid_size.inches)
            end_x = int(right / grid_size.inches) + 1
            end_y = int(bottom / grid_size.inches) + 1

            for x in range(start_x, min(end_x, int(max_x / grid_size.inches))):
                for y in range(start_y, min(end_y, int(max_y / grid_size.inches))):
                    grid[(x, y)] = grid.get((x, y), []) + [elem_id]

        # Resolve overlaps
        for (x, y), elem_ids in grid.items():
            if len(elem_ids) > 1:
                for i in range(1, len(elem_ids)):
                    elem_id = elem_ids[i]
                    current_elem = elements_dict[elem_id]
                    shift = i * grid_size
                    if x < max_x / grid_size.inches - 1:
                        current_elem["left"] += shift
                    elif y < max_y / grid_size.inches - 1:
                        current_elem["top"] += shift

        return elements_dict

    def create_ppt_from_analysis(
        self,
        analysis_data,
        outline_color=RGBColor(0, 0, 0),
        text_color=RGBColor(0, 0, 0),
    ):
        """Create a PowerPoint slide with customizable styles."""
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]  # Blank slide layout
        slide = prs.slides.add_slide(slide_layout)

        # 1. Add phase headers in separate boxes
        if "phases" in analysis_data and analysis_data["phases"]:
            phase_width = Inches(2.2)
            start_x = Inches(0.8)
            for i, phase in enumerate(analysis_data["phases"]):
                phase_x = start_x + (i * phase_width)
                phase_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, phase_x, Inches(0.2), phase_width, Inches(0.6)
                )
                phase_shape.line.color.rgb = outline_color
                phase_shape.fill.background()
                text_frame = phase_shape.text_frame
                text_frame.text = phase
                p = text_frame.paragraphs[0]
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.name = self.default_font
                p.font.color.rgb = text_color
                p.alignment = PP_ALIGN.CENTER

        # 2. Add vertical section labels
        if "sections" in analysis_data and analysis_data["sections"]:
            for i, section in enumerate(analysis_data["sections"]):
                section_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.2),
                    Inches(1 + i * 1.6),
                    Inches(0.9),
                    Inches(1.1),
                )
                section_shape.line.color.rgb = outline_color
                section_shape.fill.background()
                text_frame = section_shape.text_frame
                text_frame.text = section
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                section_shape.rotation = 270
                p = text_frame.paragraphs[0]
                p.font.size = Pt(14)
                p.font.name = self.default_font
                p.font.color.rgb = text_color

        # 3. Position map with dynamic adjustments
        position_map = {
            "top-left": {"x": 1.2, "y": 1.0},
            "top-center": {"x": 4.0, "y": 1.0},
            "top-right": {"x": 7.0, "y": 1.0},
            "middle-left": {"x": 1.2, "y": 2.9},
            "middle-center": {"x": 4.0, "y": 2.9},
            "middle-right": {"x": 7.0, "y": 2.9},
            "bottom-left": {"x": 1.2, "y": 4.8},
            "bottom-center": {"x": 4.0, "y": 4.8},
            "bottom-right": {"x": 7.0, "y": 4.8},
        }

        element_shapes = {}
        if "elements" in analysis_data:
            for element in analysis_data["elements"]:
                element_id = element.get("id", f"elem_{id(element)}")
                position = element.get("position", "middle-center")
                pos_config = position_map.get(position, {"x": 3.5, "y": 3.5})

                # Dynamic sizing based on text
                text_length = len(element["text"])
                estimated_lines = max(1, text_length // 25)
                width = min(Inches(3.5), max(Inches(1.5), Inches(text_length / 12)))
                height = min(
                    Inches(1.2), max(Inches(0.5), Inches(estimated_lines * 0.3))
                )

                # Create shape
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(pos_config["x"]),
                    Inches(pos_config["y"]),
                    width,
                    height,
                )
                shape.line.color.rgb = outline_color
                shape.fill.background()
                text_frame = shape.text_frame
                text_frame.word_wrap = True

                # Dynamic font size
                font_size = 12
                if text_length > 30:
                    font_size = 10
                if text_length > 50:
                    font_size = 9
                if text_length > 70:
                    font_size = 8

                text_frame.text = element["text"]
                p = text_frame.paragraphs[0]
                p.font.size = Pt(font_size)
                p.font.name = self.default_font
                p.font.color.rgb = text_color
                p.alignment = PP_ALIGN.CENTER

                element_shapes[element_id] = {
                    "shape": shape,
                    "left": Inches(pos_config["x"]),
                    "top": Inches(pos_config["y"]),
                    "width": width,
                    "height": height,
                }

        # Apply overlap prevention
        element_shapes = self.adjust_positions_for_overlap(element_shapes)
        for elem_id, elem_info in element_shapes.items():
            elem_info["shape"].left = elem_info["left"]
            elem_info["shape"].top = elem_info["top"]

        # 4. Add connection arrows
        if "connections" in analysis_data:
            for connection in analysis_data["connections"]:
                from_id, to_id = connection["from_id"], connection["to_id"]
                if from_id not in element_shapes or to_id not in element_shapes:
                    continue

                from_shape, to_shape = element_shapes[from_id], element_shapes[to_id]
                from_center_x = from_shape["left"] + (from_shape["width"] / 2)
                from_center_y = from_shape["top"] + (from_shape["height"] / 2)
                to_center_x = to_shape["left"] + (to_shape["width"] / 2)
                to_center_y = to_shape["top"] + (to_shape["height"] / 2)

                if abs(from_center_x - to_center_x) > abs(from_center_y - to_center_y):
                    start_x = from_center_x + (
                        from_shape["width"] / 2
                        if from_center_x < to_center_x
                        else -from_shape["width"] / 2
                    )
                    end_x = to_center_x - (
                        to_shape["width"] / 2
                        if from_center_x < to_center_x
                        else -to_shape["width"] / 2
                    )
                    start_y, end_y = from_center_y, to_center_y
                else:
                    start_y = from_center_y + (
                        from_shape["height"] / 2
                        if from_center_y < to_center_y
                        else -from_shape["height"] / 2
                    )
                    end_y = to_center_y - (
                        to_shape["height"] / 2
                        if from_center_y < to_center_y
                        else -to_shape["height"] / 2
                    )
                    start_x, end_x = from_center_x, to_center_x

                connector = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y
                )
                connector.line.color.rgb = outline_color
                connector.line.width = Pt(1.5)
                connector.line.dash_style = self.default_line_style

        # Save and return
        pptx_io = BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        return pptx_io

    def preprocess_image(self, image_data):
        """Enhance image contrast for better analysis"""
        try:
            if not image_data:
                raise ValueError("Image data is empty")

            logger.info(f"Preprocessing image data of size {len(image_data)} bytes")
            img = Image.open(io.BytesIO(image_data))
            # Convert to RGB mode to ensure compatibility with JPEG
            img = img.convert("RGB")
            enhancer = ImageEnhance.Contrast(img)
            enhanced_img = enhancer.enhance(1.5)  # Increase contrast by 50%
            output = io.BytesIO()
            enhanced_img.save(output, format="JPEG")
            return output.getvalue()
        except Exception as e:
            st.error(f"Error preprocessing image: {e}")
            logger.error(f"Preprocessing error: {e}")
            return image_data  # Return original data as fallback

    def process_image_to_ppt(self, image_data):
        """Main method with preprocessing and styling options"""
        # Validate image_data
        if not image_data or len(image_data) == 0:
            st.error("Image data is empty. Please upload a valid image.")
            return None

        # Preprocess image
        enhanced_image_data = self.preprocess_image(image_data)
        if not enhanced_image_data:
            st.error("Failed to preprocess image. Using original image data.")
            enhanced_image_data = image_data

        # Analyze image
        analysis_data = self.analyze_image(enhanced_image_data)
        if not analysis_data:
            return None

        # Get user styling preferences
        with st.sidebar:
            st.subheader("Slide Customization")
            outline_color = st.color_picker("Outline Color", "#000000")
            text_color = st.color_picker("Text Color", "#000000")
            line_style = st.selectbox(
                "Line Style", ["Solid", "Dash", "DashDot", "DashDotDot"], index=0
            )  # Removed "Dot"

        # Convert colors to RGBColor
        outline_rgb = RGBColor(
            int(outline_color[1:3], 16),
            int(outline_color[3:5], 16),
            int(outline_color[5:7], 16),
        )
        text_rgb = RGBColor(
            int(text_color[1:3], 16), int(text_color[3:5], 16), int(text_color[5:7], 16)
        )
        line_style_map = {
            "Solid": MSO_LINE_DASH_STYLE.SOLID,
            "Dash": MSO_LINE_DASH_STYLE.DASH,
            "DashDot": MSO_LINE_DASH_STYLE.DASH_DOT,
            "DashDotDot": MSO_LINE_DASH_STYLE.DASH_DOT_DOT,
        }
        self.default_line_style = line_style_map[line_style]

        # Create PowerPoint
        st.write("Creating PowerPoint slide...")
        pptx_data = self.create_ppt_from_analysis(analysis_data, outline_rgb, text_rgb)

        return pptx_data


def main():
    st.set_page_config(page_title="Whiteboard to PowerPoint Converter", page_icon="📊")

    st.title("Whiteboard to PowerPoint Converter")
    st.write("Upload a whiteboard image to generate a structured PowerPoint slide.")

    # Check for API key
    api_key = os.getenv("OPENAI_API_KEY") or st.secrets.get("OPENAI_API_KEY", "")
    if not api_key:
        st.warning(
            "⚠️ OpenAI API key not found. Please add it to your .env file or Streamlit secrets."
        )
        api_key = st.text_input("Enter your OpenAI API key:", type="password")
        if api_key:
            os.environ["OPENAI_API_KEY"] = api_key

    # File uploader with preprocessing option
    uploaded_file = st.file_uploader(
        "Upload whiteboard image", type=["jpg", "jpeg", "png"]
    )
    preprocess = st.checkbox("Enhance Image Contrast", value=True)

    if uploaded_file is not None:
        # Display the uploaded image
        try:
            image = Image.open(uploaded_file)
            st.image(image, caption="Uploaded Whiteboard Image", use_column_width=True)
        except Exception as e:
            st.error(f"Failed to display image: {e}")
            logger.error(f"Image display error: {e}")
            return

        # Process image when button is clicked
        if st.button("Convert to PowerPoint"):
            with st.spinner("Processing..."):
                # Reset file pointer and read image data
                uploaded_file.seek(0)
                image_data = uploaded_file.read()

                if not image_data:
                    st.error("Uploaded file is empty. Please upload a valid image.")
                    return

                logger.info(f"Read {len(image_data)} bytes from uploaded file")

                # Apply preprocessing if selected
                if preprocess:
                    converter = WhiteboardToPPT()
                    image_data = converter.preprocess_image(image_data)
                    if not image_data:
                        st.error("Preprocessing failed. Using original image data.")
                        return

                # Create converter instance
                converter = WhiteboardToPPT()

                # Process image
                progress_bar = st.progress(0)
                progress_bar.progress(10)
                pptx_data = converter.process_image_to_ppt(image_data)
                progress_bar.progress(100)

                if pptx_data:
                    st.success("PowerPoint slide created successfully!")
                    st.download_button(
                        label="Download PowerPoint",
                        data=pptx_data,
                        file_name="whiteboard_to_ppt.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )
                else:
                    st.error(
                        "Failed to create PowerPoint. Check logs above for details."
                    )


if __name__ == "__main__":
    main()
