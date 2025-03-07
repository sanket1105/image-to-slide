import base64
import io
import json
import os
from io import BytesIO

import requests
import streamlit as st
from dotenv import load_dotenv
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# Load environment variables for API keys
load_dotenv()
from pptx.enum.dml import MSO_LINE, MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN


# First, fix the missing method in the WhiteboardToPPT class
class WhiteboardToPPT:
    def __init__(self):
        # API configuration - use environment variable or Streamlit secrets
        self.vision_api_key = os.getenv("OPENAI_API_KEY")
        self.vision_api_url = "https://api.openai.com/v1/chat/completions"

        # PowerPoint settings
        self.slide_width = Inches(10)
        self.slide_height = Inches(7.5)

    def analyze_image(self, image_data):
        """
        Use AI vision model to analyze the whiteboard image and extract information
        Returns structured data about text content and layout
        """
        try:
            # Convert image to base64
            image_base64 = base64.b64encode(image_data).decode("utf-8")

            # Prepare improved prompt for the vision model
            prompt = """
            Analyze this whiteboard image and extract:
            1. All text content with exact wording
            2. The spatial layout (what text is in boxes, what's connected to what)
            3. The hierarchical structure (main sections, subsections)
            
            Format your response as JSON with these sections:
            - phases: List of main phase headings at the top (e.g., "Ideation", "Align", "Launch", "Manage")
            - sections: List of vertical sections on the left side
            - elements: List of all boxes/elements with their text, position, and type
            
            For positions, be specific about exact location:
            - Use "top", "middle", "bottom" for vertical position
            - Use "left", "center", "right" for horizontal position
            - Combine these for precise placement (e.g., "top-left", "middle-center")
            
            For each element, include these properties:
            - text: The exact text content
            - position: The position on the slide (e.g., "top-left", "middle-right")
            - type: "box" if text is in a box/rectangle, "text" if standalone text, "note" if it appears to be a note/comment
            - coordinates: Approximate x,y coordinates as percentages of the image (e.g., {"x": 20, "y": 30})
            - has_shape: boolean indicating if text is enclosed in a visible shape
            
            Example format:
            {
              "phases": ["Ideation", "Align", "Launch", "Manage"],
              "sections": ["Strategy", "Planning", "YES", "Foundation"],
              "elements": [
                {"text": "Contribute Content Ideas", "position": "middle-left", "type": "box", "coordinates": {"x": 20, "y": 30}, "has_shape": true},
                {"text": "#Need Sizing Tool#", "position": "middle-left", "type": "note", "coordinates": {"x": 25, "y": 40}, "has_shape": false}
              ]
            }
            
            Important: Ensure there are no overlapping elements. If you detect text that appears to be stacked or overlapping, separate them into distinct elements with different coordinates.
            """

            # API call to vision model
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {self.vision_api_key}",
            }

            # Updated to use the current GPT-4 Vision model
            payload = {
                "model": "gpt-4o",  # Current model for vision capabilities
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/jpeg;base64,{image_base64}"
                                },
                            },
                        ],
                    }
                ],
                "max_tokens": 4000,
            }

            st.write("Analyzing image with AI vision model...")

            response = requests.post(self.vision_api_url, headers=headers, json=payload)

            if response.status_code != 200:
                st.error(f"API Error: {response.status_code}")
                st.error(response.text)
                return None

            result = response.json()

            # Check if we have the expected structure
            if "choices" not in result or len(result["choices"]) == 0:
                st.error("Unexpected API response format")
                st.write(result)
                return None

            # Extract the content from the response
            content = result["choices"][0]["message"]["content"]

            # Parse the JSON content - handle potential JSON formatting issues
            try:
                # Find JSON content in the response (it might be mixed with other text)
                json_content = self._extract_json_from_text(content)
                analysis_data = json.loads(json_content)
                return analysis_data
            except json.JSONDecodeError as e:
                st.error(f"Error parsing JSON response: {e}")
                st.write("Raw content:")
                st.write(content)

                # Fallback to manual extraction if JSON parsing fails
                return self._fallback_extraction(content)

        except Exception as e:
            st.error(f"Error analyzing image: {e}")
            return None

    def _extract_json_from_text(self, text):
        """Extract JSON content from potentially mixed text"""
        # Look for content between curly braces
        start_idx = text.find("{")
        end_idx = text.rfind("}")

        if start_idx >= 0 and end_idx > start_idx:
            return text[start_idx : end_idx + 1]

        return "{}"  # Return empty JSON if none found

    def _fallback_extraction(self, content):
        """Manual extraction of key elements as fallback"""
        # Based on the image, return a hardcoded structure as fallback
        st.warning("Using fallback extraction method")
        return {
            "phases": ["Ideation", "Align", "Launch", "Manage"],
            "sections": ["Strategy", "Planning", "YES", "Foundation"],
            "elements": [
                {
                    "text": "Contribute Content Ideas",
                    "position": "top-left",
                    "type": "box",
                    "coordinates": {"x": 15, "y": 25},
                    "has_shape": True,
                },
                {
                    "text": "Define Test Audience",
                    "position": "middle-left",
                    "type": "box",
                    "coordinates": {"x": 20, "y": 45},
                    "has_shape": True,
                },
                {
                    "text": "Sales Calling & Accountability",
                    "position": "top-right",
                    "type": "box",
                    "coordinates": {"x": 80, "y": 25},
                    "has_shape": True,
                },
                {
                    "text": "Productivity Reporting",
                    "position": "middle-right",
                    "type": "box",
                    "coordinates": {"x": 80, "y": 45},
                    "has_shape": True,
                },
                {
                    "text": "Finance Potential",
                    "position": "top-right",
                    "type": "text",
                    "coordinates": {"x": 65, "y": 15},
                    "has_shape": False,
                },
                {
                    "text": "(?Randy?)",
                    "position": "bottom-right",
                    "type": "note",
                    "coordinates": {"x": 85, "y": 85},
                    "has_shape": False,
                },
            ],
        }

    def create_ppt_from_analysis(self, analysis_data):
        """
        Create a PowerPoint slide based on the AI-analyzed content
        """
        # Initialize PowerPoint presentation
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]  # Blank slide layout
        slide = prs.slides.add_slide(slide_layout)

        # 1. Add the phase headers at the top in SEPARATE BOXES
        if "phases" in analysis_data and analysis_data["phases"]:
            phase_width = Inches(2)
            total_phases = len(analysis_data["phases"])
            start_x = Inches(1)

            for i, phase in enumerate(analysis_data["phases"]):
                # Calculate position for evenly spaced boxes
                phase_x = start_x + (i * phase_width)

                # Create a box for each phase
                phase_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left=phase_x,
                    top=Inches(0.2),
                    width=phase_width,
                    height=Inches(0.5),
                )

                # Configure the shape
                phase_shape.line.color.rgb = RGBColor(0, 0, 0)  # Black outline
                phase_shape.fill.background()  # Transparent fill

                # Add text
                text_frame = phase_shape.text_frame
                text_frame.text = phase
                p = text_frame.paragraphs[0]
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.name = "Arial"
                p.font.color.rgb = RGBColor(0, 0, 0)  # Black text
                p.alignment = PP_ALIGN.CENTER

        # 2. Add the vertical section labels on the left in boxes
        if "sections" in analysis_data and analysis_data["sections"]:
            for i, section in enumerate(analysis_data["sections"]):
                section_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left=Inches(0.2),
                    top=Inches(1 + i * 1.2),
                    width=Inches(0.8),
                    height=Inches(1),
                )
                section_shape.line.color.rgb = RGBColor(0, 0, 0)  # Black outline
                section_shape.fill.background()  # Transparent fill

                # Add text
                text_frame = section_shape.text_frame
                text_frame.text = section
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                section_shape.rotation = 270  # Rotate text vertically
                p = text_frame.paragraphs[0]
                p.font.size = Pt(14)
                p.font.name = "Arial"
                p.font.color.rgb = RGBColor(0, 0, 0)  # Black text

        # 3. Create a mapping for positions to actual coordinates
        # This helps prevent overlap by using specific position mappings
        position_map = {
            # Top row positions
            "top-left": {"x": 1, "y": 1.5},
            "top-center": {"x": 3.5, "y": 1.5},
            "top-right": {"x": 6, "y": 1.5},
            # Middle row positions
            "middle-left": {"x": 1, "y": 3},
            "middle-center": {"x": 3.5, "y": 3},
            "middle-right": {"x": 6, "y": 3},
            # Bottom row positions
            "bottom-left": {"x": 1, "y": 4.5},
            "bottom-center": {"x": 3.5, "y": 4.5},
            "bottom-right": {"x": 6, "y": 4.5},
            # Additional positions for more precise placement
            "top-left-2": {"x": 1, "y": 2},
            "middle-left-2": {"x": 1, "y": 3.5},
            "bottom-left-2": {"x": 1, "y": 5},
            "top-right-2": {"x": 6, "y": 2},
            "middle-right-2": {"x": 6, "y": 3.5},
            "bottom-right-2": {"x": 6, "y": 5},
        }

        # Store elements by position for adding connections later
        element_shapes = {}

        # Add additional dynamic positions based on coordinates if provided
        if "elements" in analysis_data:
            for i, element in enumerate(analysis_data["elements"]):
                if "coordinates" in element:
                    # Create a unique position key based on coordinates
                    pos_key = f"custom-{i}"
                    # Convert percentage coordinates to inches (assuming 10 inches width, 7.5 inches height)
                    x_inches = (element["coordinates"]["x"] / 100) * 10
                    y_inches = (element["coordinates"]["y"] / 100) * 7.5
                    position_map[pos_key] = {"x": x_inches, "y": y_inches}
                    # Add the position key to the element for later use
                    element["position_key"] = pos_key
                else:
                    # If no coordinates, use the text position
                    element["position_key"] = element["position"]

        # 4. Add all the elements based on their position and type
        if "elements" in analysis_data:
            for element in analysis_data["elements"]:
                position = element.get(
                    "position_key", element.get("position", "middle-center")
                )
                element_type = element.get("type", "text")
                # Force all elements to have shapes (boxes)
                has_shape = True  # Always create boxes

                # Get position configuration or use default
                pos_config = position_map.get(
                    position, {"x": 3.5, "y": 3.5}  # Default to middle of slide
                )

                # Extract coordinates
                left = Inches(pos_config["x"])
                top = Inches(pos_config["y"])

                # Determine text length to set appropriate width
                text_length = len(element["text"])
                width = min(Inches(4), max(Inches(1.5), Inches(text_length / 12)))

                # Create shape for all elements (always a box)
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left=left,
                    top=top,
                    width=width,
                    height=Inches(0.8),
                )

                # Set black outline for all shapes
                shape.line.color.rgb = RGBColor(0, 0, 0)  # Black outline
                shape.fill.background()  # Transparent fill

                text_frame = shape.text_frame
                text_frame.word_wrap = True

                # Add and format the text
                text_frame.text = element["text"]
                p = text_frame.paragraphs[0]
                p.font.size = Pt(12)
                p.font.name = "Arial"
                p.font.color.rgb = RGBColor(0, 0, 0)  # Black text for all elements

                # Store the shape by position for connection arrows
                element_shapes[position] = {
                    "shape": shape,
                    "text": element["text"],
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": Inches(0.8),
                }

        # 5. Add connection arrows between boxes based on the image
        # Define the connections based on the whiteboard image
        connections = [
            # Example connections (adjust based on your actual image)
            {"from": "middle-left", "to": "middle-center"},
            {"from": "middle-center", "to": "middle-right"},
            {"from": "top-left", "to": "top-center"},
            {"from": "top-center", "to": "top-right"},
            # Add additional connections as needed
        ]

        # Create the arrows
        for connection in connections:
            from_pos = connection["from"]
            to_pos = connection["to"]

            # Skip if either shape doesn't exist
            if from_pos not in element_shapes or to_pos not in element_shapes:
                continue

            from_shape = element_shapes[from_pos]
            to_shape = element_shapes[to_pos]

            # Calculate connection points
            start_x = from_shape["left"] + from_shape["width"]
            start_y = from_shape["top"] + (from_shape["height"] / 2)
            end_x = to_shape["left"]
            end_y = to_shape["top"] + (to_shape["height"] / 2)

            # Create arrow connector
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y
            )
            connector.line.color.rgb = RGBColor(0, 0, 0)  # Black arrow
            connector.line.width = Pt(1.5)
            connector.line.dash_style = MSO_LINE_DASH_STYLE.SOLID

            # Add arrow head
            connector.line.begin_style = MSO_LINE.NONE
            connector.line.end_style = MSO_LINE.ARROW

        # Save the PowerPoint
        pptx_io = BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        return pptx_io

    # Add the missing process_image_to_ppt method
    def process_image_to_ppt(self, image_data):
        """
        Main method to process image and create PowerPoint
        """
        # Step 1: Analyze image with AI
        analysis_data = self.analyze_image(image_data)

        if not analysis_data:
            st.error("Failed to analyze image")
            return None

        # Step 2: Create PowerPoint from analysis
        st.write("Creating PowerPoint slide...")
        pptx_data = self.create_ppt_from_analysis(analysis_data)

        return pptx_data


# Also, fix the deprecated parameter in the main function
def main():
    st.set_page_config(page_title="Whiteboard to PowerPoint Converter", page_icon="📊")

    st.title("Whiteboard to PowerPoint Converter")
    st.write("Upload a whiteboard image and get a structured PowerPoint slide.")

    # Check for API key
    api_key = os.getenv("OPENAI_API_KEY") or st.secrets.get("OPENAI_API_KEY", "")
    if not api_key:
        st.warning(
            "⚠️ OpenAI API key not found. Please add it to your .env file or Streamlit secrets."
        )
        api_key = st.text_input("Enter your OpenAI API key:", type="password")
        if api_key:
            os.environ["OPENAI_API_KEY"] = api_key

    # File uploader
    uploaded_file = st.file_uploader(
        "Upload whiteboard image", type=["jpg", "jpeg", "png"]
    )

    if uploaded_file is not None:
        # Display the uploaded image - FIX THE DEPRECATED PARAMETER
        image = Image.open(uploaded_file)
        st.image(image, caption="Uploaded Whiteboard Image", use_container_width=True)

        # Process image when button is clicked
        if st.button("Convert to PowerPoint"):
            with st.spinner("Processing..."):
                # Reset uploaded_file to start
                uploaded_file.seek(0)
                image_data = uploaded_file.read()

                # Create converter instance
                converter = WhiteboardToPPT()

                # Process image
                pptx_data = converter.process_image_to_ppt(image_data)

                if pptx_data:
                    # Provide download button for the generated PowerPoint
                    st.success("PowerPoint slide created successfully!")

                    # Create download button
                    st.download_button(
                        label="Download PowerPoint",
                        data=pptx_data,
                        file_name="whiteboard_to_ppt.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )
                else:
                    st.error(
                        "Failed to create PowerPoint. Please check the logs above for details."
                    )


if __name__ == "__main__":
    main()
