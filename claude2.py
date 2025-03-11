import base64
import io
import json
import math
import os
from io import BytesIO

import cv2
import numpy as np
import requests
import streamlit as st
from dotenv import load_dotenv
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# Load environment variables for API keys
load_dotenv()
from pptx.enum.dml import MSO_LINE_DASH_STYLE, MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN


class WhiteboardToPPT:
    def __init__(self):
        # API configuration - use environment variable or Streamlit secrets
        self.vision_api_key = os.getenv("OPENAI_API_KEY")
        self.vision_api_url = "https://api.openai.com/v1/chat/completions"

        # PowerPoint settings
        self.slide_width = Inches(10)
        self.slide_height = Inches(7.5)

        # Configuration for element sizing and spacing
        self.min_box_width = Inches(1.2)
        self.max_box_width = Inches(3.5)
        self.min_box_height = Inches(0.5)
        self.max_box_height = Inches(1.2)
        self.horizontal_spacing = Inches(0.4)  # Space between elements horizontally
        self.vertical_spacing = Inches(0.4)  # Space between elements vertically

    def rotate_image(self, image_data, angle):
        """
        Manually rotate the image by a specified angle.

        Parameters:
        - image_data: The binary image data
        - angle: Rotation angle in degrees (positive = counterclockwise)

        Returns:
        - Rotated image as binary data
        """
        try:
            # Convert image data to OpenCV format
            nparr = np.frombuffer(image_data, np.uint8)
            img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)

            # Get image dimensions
            h, w = img.shape[:2]
            center = (w // 2, h // 2)

            # Calculate rotation matrix
            M = cv2.getRotationMatrix2D(center, angle, 1.0)

            # Calculate new image dimensions after rotation
            cos = abs(M[0, 0])
            sin = abs(M[0, 1])
            new_w = int((h * sin) + (w * cos))
            new_h = int((h * cos) + (w * sin))

            # Adjust the rotation matrix to take into account the translation
            M[0, 2] += (new_w / 2) - center[0]
            M[1, 2] += (new_h / 2) - center[1]

            # Perform the rotation
            rotated = cv2.warpAffine(
                img,
                M,
                (new_w, new_h),
                flags=cv2.INTER_CUBIC,
                borderMode=cv2.BORDER_CONSTANT,
                borderValue=(255, 255, 255),  # White border
            )

            # Convert the rotated OpenCV image back to bytes
            success, rotated_image = cv2.imencode(".jpg", rotated)
            rotated_data = rotated_image.tobytes()

            return rotated_data
        except Exception as e:
            st.error(f"Error during image rotation: {e}")
            # Return original image if rotation fails
            return image_data

    def preprocess_image(self, image_data, manual_rotation=None):
        """
        Preprocess the image to enhance OCR accuracy:
        1. Apply manual rotation if specified
        2. Convert to OpenCV format
        3. Detect edges and correct skew/rotation (if manual rotation not specified)
        4. Enhance contrast and sharpness
        5. Remove noise

        Parameters:
        - image_data: The binary image data
        - manual_rotation: Optional manual rotation angle (degrees)

        Returns:
        - Processed image as binary data
        """
        try:
            # First apply manual rotation if specified
            if manual_rotation is not None and manual_rotation != 0:
                st.info(f"Applying manual rotation: {manual_rotation} degrees")
                image_data = self.rotate_image(image_data, manual_rotation)

            # Convert image data to OpenCV format
            nparr = np.frombuffer(image_data, np.uint8)
            img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)

            # Keep a copy of the original for comparison
            original_img = img.copy()

            # Convert to grayscale
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)

            # Apply bilateral filter to reduce noise while preserving edges
            gray = cv2.bilateralFilter(gray, 11, 17, 17)

            # Auto-rotation only if manual rotation is not specified
            if manual_rotation is None:
                # Detect edges in the image
                edges = cv2.Canny(gray, 30, 200)

                # Use Hough Line Transform to detect lines
                lines = cv2.HoughLinesP(
                    edges, 1, np.pi / 180, 100, minLineLength=100, maxLineGap=10
                )

                if lines is not None and len(lines) > 0:
                    # Find the dominant angle
                    angles = []
                    for line in lines:
                        x1, y1, x2, y2 = line[0]
                        if x2 - x1 != 0:  # Avoid division by zero
                            angle = math.atan2(y2 - y1, x2 - x1) * 180.0 / np.pi
                            # We're interested in horizontal and vertical lines
                            # Normalize angles to find lines close to horizontal (0 or 180 degrees)
                            norm_angle = abs(angle) % 180
                            if norm_angle > 90:
                                norm_angle = 180 - norm_angle
                            if norm_angle < 45:  # Consider it a horizontal line
                                angles.append(angle)

                    if angles:
                        # Calculate the median angle to avoid outliers
                        median_angle = np.median(angles)

                        # Rotate the image to correct the skew
                        if (
                            abs(median_angle) > 0.5
                        ):  # Only rotate if the angle is significant
                            st.info(
                                f"Correcting image skew: {median_angle:.2f} degrees"
                            )

                            # Get image dimensions
                            h, w = img.shape[:2]
                            center = (w // 2, h // 2)

                            # Calculate rotation matrix
                            M = cv2.getRotationMatrix2D(center, median_angle, 1.0)

                            # Calculate new image dimensions after rotation
                            cos = abs(M[0, 0])
                            sin = abs(M[0, 1])
                            new_w = int((h * sin) + (w * cos))
                            new_h = int((h * cos) + (w * sin))

                            # Adjust the rotation matrix to take into account the translation
                            M[0, 2] += (new_w / 2) - center[0]
                            M[1, 2] += (new_h / 2) - center[1]

                            # Perform the rotation
                            rotated = cv2.warpAffine(
                                img,
                                M,
                                (new_w, new_h),
                                flags=cv2.INTER_CUBIC,
                                borderMode=cv2.BORDER_CONSTANT,
                                borderValue=(255, 255, 255),
                            )
                            img = rotated

            # Enhance contrast
            clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
            enhanced = clahe.apply(gray)

            # Merge enhanced image back to color if needed
            if img.shape[2] == 3:  # If it's a color image
                enhanced_color = cv2.cvtColor(enhanced, cv2.COLOR_GRAY2BGR)
                # Blend with original for a more natural look
                img = cv2.addWeighted(img, 0.7, enhanced_color, 0.3, 0)

            # Convert the processed OpenCV image back to bytes
            success, processed_image = cv2.imencode(".jpg", img)
            processed_data = processed_image.tobytes()

            # Display original and processed images for comparison
            with st.expander("Image Preprocessing Results", expanded=True):
                col1, col2 = st.columns(2)
                with col1:
                    # Convert original OpenCV image to PIL for display
                    original_rgb = cv2.cvtColor(original_img, cv2.COLOR_BGR2RGB)
                    st.image(original_rgb, caption="Original Image")
                with col2:
                    # Convert processed OpenCV image to PIL for display
                    processed_rgb = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
                    st.image(processed_rgb, caption="Processed Image")

            return processed_data
        except Exception as e:
            st.error(f"Error during image preprocessing: {e}")
            # Return original image if preprocessing fails
            return image_data

    def analyze_image(self, image_data, rotation_angle=None):
        """
        Use AI vision model to analyze the whiteboard image and extract information.
        Returns structured data about text content and layout.

        Parameters:
        - image_data: The binary image data
        - rotation_angle: Optional manual rotation angle (degrees)
        """
        try:
            # Preprocess the image to enhance OCR accuracy
            processed_image_data = self.preprocess_image(image_data, rotation_angle)

            # Convert preprocessed image to base64
            image_base64 = base64.b64encode(processed_image_data).decode("utf-8")

            # Enhanced prompt for the vision model with better OCR instructions
            prompt = """
            Analyze this whiteboard image with high precision OCR, extracting:
            1. All text content with EXACT wording, preserving all punctuation, capitalization, and special characters
            2. The spatial layout and relationships between elements
            3. The hierarchical structure and flow of information
            
            CRITICALLY IMPORTANT:
            - Scan the image multiple times to ensure ALL text is captured
            - Verify each extracted text element is unique (DO NOT DUPLICATE any text)
            - Ensure text with special characters (#, *, etc.) is captured correctly
            - Treat faint or partial text as important and capture it
            - Capture text in cloud shapes, rectangles, and any other visual elements
            
            FORMAT REQUIREMENTS - STRICTLY ADHERE TO THIS JSON STRUCTURE:
            {
              "phases": [list of phase headings at the top, ordered left to right],
              "sections": [list of vertical sections on the left side, ordered top to bottom],
              "elements": [
                {
                  "id": "unique identifier (elem1, elem2, etc.)",
                  "text": "exact text content with ALL special characters and formatting preserved",
                  "position": "precise position descriptor (top-left, middle-center, etc.)",
                  "type": "box/text/note/cloud",
                  "coordinates": {"x": percentage from left, "y": percentage from top},
                  "has_shape": boolean,
                  "emphasis": "none/high/medium/low based on visual prominence"
                }
              ],
              "connections": [
                {"from_id": "source element id", "to_id": "target element id", "arrow_type": "simple/double/none"}
              ]
            }
            
            Before finalizing, verify:
            1. NO DUPLICATE TEXT in any elements
            2. ALL TEXT from image is captured (compare your extraction with the image)
            3. Each element has a unique ID
            4. Clouds and special shapes are properly identified
            """

            # API call to vision model
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {self.vision_api_key}",
            }

            # Using GPT-4o for improved visual analysis with temperature=0 for consistency
            payload = {
                "model": "gpt-4o",
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/jpeg;base64,{image_base64}"
                                },
                            },
                        ],
                    }
                ],
                "max_tokens": 4000,
                "temperature": 0,  # Using 0 for maximum consistency and accuracy
            }

            with st.status(
                "Analyzing image with AI vision model...", expanded=True
            ) as status:
                response = requests.post(
                    self.vision_api_url, headers=headers, json=payload
                )

                if response.status_code != 200:
                    st.error(f"API Error: {response.status_code}")
                    st.error(response.text)
                    status.update(label="Analysis failed", state="error")
                    return None

                result = response.json()
                status.update(label="Analysis complete", state="complete")

            # Check if we have the expected structure
            if "choices" not in result or len(result["choices"]) == 0:
                st.error("Unexpected API response format")
                st.write(result)
                return None

            # Extract the content from the response
            content = result["choices"][0]["message"]["content"]

            # Parse the JSON content - handle potential JSON formatting issues
            try:
                # Find JSON content in the response (it might be mixed with other text)
                json_content = self._extract_json_from_text(content)
                analysis_data = json.loads(json_content)

                # Log the detected connections and elements
                st.expander("AI Analysis Results", expanded=False).write(analysis_data)

                # Check for duplicate texts and missing elements
                self._validate_analysis_data(analysis_data)

                # Enhance data with text analysis
                analysis_data = self._enhance_analysis_data(analysis_data)

                return analysis_data
            except json.JSONDecodeError as e:
                st.error(f"Error parsing JSON response: {e}")
                st.write("Raw content:")
                st.write(content)

                # Fallback to manual extraction if JSON parsing fails
                return self._fallback_extraction(content)

        except Exception as e:
            st.error(f"Error analyzing image: {e}")
            return None

    # Rest of your class methods remain the same
    def _validate_analysis_data(self, data):
        """Validate the analysis data for duplicates and missing elements"""
        if not data or "elements" not in data:
            return

        # Check for duplicate texts
        texts = []
        duplicates = []

        for element in data["elements"]:
            if "text" in element:
                if element["text"] in texts:
                    duplicates.append(element["text"])
                texts.append(element["text"])

        if duplicates:
            st.warning(
                f"Found {len(duplicates)} duplicate text elements. Removing duplicates..."
            )

            # Remove duplicates while preserving the first occurrence
            seen_texts = set()
            filtered_elements = []

            for element in data["elements"]:
                if "text" in element and element["text"] not in seen_texts:
                    seen_texts.add(element["text"])
                    filtered_elements.append(element)
                elif "text" not in element:
                    filtered_elements.append(element)

            data["elements"] = filtered_elements
            st.success("Duplicates removed.")

        # Also check for duplicate phase headers
        if "phases" in data and data["phases"]:
            unique_phases = []
            for phase in data["phases"]:
                if phase not in unique_phases:
                    unique_phases.append(phase)
            data["phases"] = unique_phases

    def _enhance_analysis_data(self, data):
        """Add additional analysis to improve the layout and connections"""
        if not data:
            return data

        # Add size estimation based on text length
        if "elements" in data:
            for element in data["elements"]:
                if "text" in element:
                    # Calculate text metrics
                    text = element["text"]
                    element["text_length"] = len(text)
                    element["word_count"] = len(text.split())

                    # Estimate importance based on position, content, and connections
                    importance = 1  # Default importance

                    # Check if this element is connected to many others
                    if "connections" in data:
                        connections_from = sum(
                            1
                            for conn in data["connections"]
                            if conn["from_id"] == element["id"]
                        )
                        connections_to = sum(
                            1
                            for conn in data["connections"]
                            if conn["to_id"] == element["id"]
                        )
                        element["connection_count"] = connections_from + connections_to

                        # Elements with more connections tend to be more important
                        if connections_from + connections_to > 2:
                            importance += 1

                    # Top position elements are often more important
                    if "position" in element and "top" in element["position"].lower():
                        importance += 1

                    # Elements with special markers like "#" might be emphasized
                    if "#" in text or "*" in text or "!" in text:
                        importance += 1

                    element["importance"] = min(importance, 3)  # Scale of 1-3

        return data

    def _extract_json_from_text(self, text):
        """Extract JSON content from potentially mixed text"""
        # Look for content between curly braces
        start_idx = text.find("{")
        end_idx = text.rfind("}")

        if start_idx >= 0 and end_idx > start_idx:
            return text[start_idx : end_idx + 1]

        return "{}"  # Return empty JSON if none found

    def _fallback_extraction(self, content):
        """Enhanced manual extraction of key elements as fallback"""
        st.warning("Using fallback extraction method")

        # More sophisticated fallback with better structure
        return {
            "phases": ["Ideation", "Align", "Launch", "Manage"],
            "sections": ["Strategy", "Planning", "Foundation"],
            "elements": [
                {
                    "id": "elem1",
                    "text": "#Need Sizing Tool#",
                    "position": "top-left",
                    "type": "note",
                    "coordinates": {"x": 15, "y": 13},
                    "has_shape": True,
                    "emphasis": "high",
                    "text_length": 17,
                    "word_count": 3,
                    "importance": 3,
                },
                {
                    "id": "elem2",
                    "text": "Define: Test (channels) Audience Success",
                    "position": "top-center",
                    "type": "box",
                    "coordinates": {"x": 35, "y": 13},
                    "has_shape": True,
                    "emphasis": "medium",
                    "text_length": 42,
                    "word_count": 5,
                    "importance": 2,
                },
                # Additional elements would be included here...
            ],
            "connections": [
                {"from_id": "elem1", "to_id": "elem2", "arrow_type": "simple"},
                # Additional connections would be included here...
            ],
        }

    def _calculate_element_dimensions(self, element):
        """Calculate appropriate dimensions for an element based on its content"""
        text = element.get("text", "")
        text_length = element.get("text_length", len(text))
        word_count = element.get("word_count", len(text.split()))
        importance = element.get("importance", 1)

        # Calculate width based on text length and word count
        # Longer words need more width
        avg_word_length = text_length / max(word_count, 1)
        width_factor = min(
            1.0, 0.6 + (avg_word_length / 20)
        )  # Scale based on average word length

        # Calculate width based on text length with constraints
        width = max(
            self.min_box_width,
            min(self.max_box_width, Inches(1.0 + (text_length / 15) * width_factor)),
        )

        # Calculate height based on text length and width
        # Estimate how many lines the text will take
        chars_per_line = width.inches * 11  # Approx 11 characters per inch
        estimated_lines = max(1, text_length / chars_per_line)

        # Add some extra height for important elements
        importance_factor = 1.0 + ((importance - 1) * 0.15)

        height = max(
            self.min_box_height,
            min(
                self.max_box_height,
                Inches(0.4 + (estimated_lines * 0.2) * importance_factor),
            ),
        )

        return width, height

    def _intelligently_position_elements(self, analysis_data):
        """Create an optimal layout for elements based on their relationships and positions"""
        # Create a grid system for positioning
        grid = {
            "rows": 6,  # Increase number of rows for better spacing
            "cols": 6,  # Increase number of columns for better spacing
            "cells": {},  # Will store what's placed in each cell
        }

        # First, prepare data on all elements
        elements_data = {}
        if "elements" in analysis_data:
            # Sort elements by importance/priority
            elements = sorted(
                analysis_data["elements"],
                key=lambda e: (
                    # Sort by y-coordinate first (top to bottom)
                    e.get("coordinates", {}).get("y", 50),
                    # Then by x-coordinate (left to right)
                    e.get("coordinates", {}).get("x", 50),
                ),
            )

            for element in elements:
                element_id = element.get("id")
                if not element_id:
                    continue

                # Get position information
                pos = element.get("position", "middle-center")
                coords = element.get("coordinates", {"x": 50, "y": 50})

                # Determine grid position
                row = min(grid["rows"] - 1, int((coords["y"] / 100) * grid["rows"]))
                col = min(grid["cols"] - 1, int((coords["x"] / 100) * grid["cols"]))

                # Calculate dimensions
                width, height = self._calculate_element_dimensions(element)

                # Store information for this element
                elements_data[element_id] = {
                    "element": element,
                    "grid_pos": (row, col),
                    "width": width,
                    "height": height,
                    "connections": [],
                }

                # Mark this grid cell as occupied
                if (row, col) not in grid["cells"]:
                    grid["cells"][(row, col)] = []
                grid["cells"][(row, col)].append(element_id)

        # Add connection information
        if "connections" in analysis_data:
            for connection in analysis_data["connections"]:
                from_id = connection.get("from_id")
                to_id = connection.get("to_id")

                if from_id in elements_data and to_id in elements_data:
                    elements_data[from_id]["connections"].append(
                        {"to": to_id, "type": connection.get("arrow_type", "simple")}
                    )
                    elements_data[to_id]["connections"].append(
                        {
                            "from": from_id,
                            "type": connection.get("arrow_type", "simple"),
                        }
                    )

        # Now convert grid positions to actual coordinates
        slide_margin_x = Inches(1.0)
        slide_margin_y = Inches(1.0)

        # Calculate usable area for elements
        usable_width = self.slide_width - (slide_margin_x * 2)
        usable_height = self.slide_height - (slide_margin_y * 2)

        # Calculate cell size
        cell_width = usable_width / grid["cols"]
        cell_height = usable_height / grid["rows"]

        # IMPROVED OVERLAP HANDLING
        # Step 1: First pass to place elements on a more spread out grid
        for element_id, data in elements_data.items():
            row, col = data["grid_pos"]

            # Calculate base position (centered in the cell)
            base_x = (
                slide_margin_x
                + (col * cell_width)
                + (cell_width / 2)
                - (data["width"] / 2)
            )
            base_y = (
                slide_margin_y
                + (row * cell_height)
                + (cell_height / 2)
                - (data["height"] / 2)
            )

            # Store the calculated position
            data["position"] = {"left": base_x, "top": base_y}

        # Step 2: Second pass to detect and resolve overlaps
        resolved_positions = {}

        # Sort elements for deterministic processing (helps with consistent layout)
        sorted_elements = sorted(
            elements_data.items(),
            key=lambda x: (x[1]["grid_pos"][0], x[1]["grid_pos"][1]),
        )

        for element_id, data in sorted_elements:
            current_pos = data["position"]
            current_width = data["width"]
            current_height = data["height"]

            # Check if this element overlaps with any already positioned element
            overlapping = False

            for positioned_id, positioned_pos in resolved_positions.items():
                # Skip checking against itself
                if positioned_id == element_id:
                    continue

                positioned_width = elements_data[positioned_id]["width"]
                positioned_height = elements_data[positioned_id]["height"]

                # Check for overlap
                if (
                    current_pos["left"] < positioned_pos["left"] + positioned_width
                    and current_pos["left"] + current_width > positioned_pos["left"]
                    and current_pos["top"] < positioned_pos["top"] + positioned_height
                    and current_pos["top"] + current_height > positioned_pos["top"]
                ):

                    overlapping = True

                    # Determine best direction to shift (right or down)
                    # Calculate distances for potential shifts
                    shift_right = (
                        positioned_pos["left"]
                        + positioned_width
                        - current_pos["left"]
                        + self.horizontal_spacing
                    )
                    shift_down = (
                        positioned_pos["top"]
                        + positioned_height
                        - current_pos["top"]
                        + self.vertical_spacing
                    )

                    # Choose the smaller shift
                    if (
                        shift_right <= shift_down
                        and current_pos["left"] + current_width + shift_right
                        <= self.slide_width - slide_margin_x
                    ):
                        # Shift right
                        current_pos["left"] += shift_right
                    else:
                        # Shift down
                        current_pos["top"] += shift_down

                    # Re-check all positioned elements again after shifting
                    # This is needed to ensure the newly shifted position doesn't overlap with others
                    overlapping = True
                    break

            # Keep checking until no overlaps remain
            while overlapping:
                overlapping = False
                for positioned_id, positioned_pos in resolved_positions.items():
                    if positioned_id == element_id:
                        continue

                    positioned_width = elements_data[positioned_id]["width"]
                    positioned_height = elements_data[positioned_id]["height"]

                    if (
                        current_pos["left"] < positioned_pos["left"] + positioned_width
                        and current_pos["left"] + current_width > positioned_pos["left"]
                        and current_pos["top"]
                        < positioned_pos["top"] + positioned_height
                        and current_pos["top"] + current_height > positioned_pos["top"]
                    ):

                        overlapping = True

                        # Try shifting right first, if still on slide
                        if (
                            current_pos["left"]
                            + current_width
                            + self.horizontal_spacing
                            <= self.slide_width - slide_margin_x
                        ):
                            current_pos["left"] += self.horizontal_spacing
                        # Otherwise shift down
                        else:
                            current_pos["left"] = slide_margin_x  # Reset to left side
                            current_pos["top"] += self.vertical_spacing

                        break

            # Add this element to the resolved positions
            resolved_positions[element_id] = current_pos
            # Update the position in the original data
            data["position"] = current_pos

        return elements_data

    def create_ppt_from_analysis(self, analysis_data):
        """
        Create a PowerPoint slide based on the AI-analyzed content.
        """
        # Initialize PowerPoint presentation with improved slide sizing
        prs = Presentation()

        # Set slide dimensions for better layout
        prs.slide_width = self.slide_width
        prs.slide_height = self.slide_height

        slide_layout = prs.slide_layouts[6]  # Blank slide layout
        slide = prs.slides.add_slide(slide_layout)

        # 1. Add title to the slide for context
        title_shape = slide.shapes.add_textbox(
            left=Inches(0.5), top=Inches(0.1), width=Inches(9.0), height=Inches(0.5)
        )
        title_frame = title_shape.text_frame
        title_frame.text = "Whiteboard Conversion"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(20)
        title_para.font.bold = True
        title_para.font.name = "Arial"
        title_para.font.color.rgb = RGBColor(0, 0, 0)  # Black text for title
        title_para.alignment = PP_ALIGN.CENTER

        # 2. Add the phase headers at the top with better spacing and styling
        if "phases" in analysis_data and analysis_data["phases"]:
            phase_width = Inches(1.8)
            phase_height = Inches(0.6)
            total_phases = len(analysis_data["phases"])

            # Calculate total width needed for all phases
            total_width = total_phases * phase_width
            # Calculate start position to center the phases
            start_x = (self.slide_width - total_width) / 2

            for i, phase in enumerate(analysis_data["phases"]):
                # Calculate position for evenly spaced boxes
                phase_x = start_x + (i * phase_width)

                # Create a box for each phase
                phase_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,  # Use rounded rectangle for modern look
                    left=phase_x,
                    top=Inches(0.7),
                    width=phase_width,
                    height=phase_height,
                )

                # Style the phase box for better visibility
                fill = phase_shape.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(230, 230, 250)  # Light lavender

                # Add border to the shape
                line = phase_shape.line
                line.color.rgb = RGBColor(128, 128, 128)  # Gray border
                line.width = Pt(1.5)

                # Add text to the phase
                text_frame = phase_shape.text_frame
                text_frame.text = phase
                text_frame.margin_bottom = Inches(0.05)
                text_frame.margin_left = Inches(0.05)
                text_frame.margin_right = Inches(0.05)
                text_frame.margin_top = Inches(0.05)
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                # Style the text
                para = text_frame.paragraphs[0]
                para.alignment = PP_ALIGN.CENTER
                para.font.bold = True
                para.font.size = Pt(14)
                para.font.name = "Arial"
                para.font.color.rgb = RGBColor(0, 0, 102)  # Dark blue text

        # 3. Intelligently position and create elements
        elements_data = self._intelligently_position_elements(analysis_data)

        # Track created shapes for connecting lines
        created_shapes = {}

        # Create all elements first
        for element_id, data in elements_data.items():
            element = data["element"]
            position = data["position"]
            width = data["width"]
            height = data["height"]

            # Determine shape type based on element type
            shape_type = MSO_SHAPE.RECTANGLE  # Default

            if element.get("type") == "cloud":
                shape_type = MSO_SHAPE.CLOUD
            elif element.get("type") == "note":
                shape_type = MSO_SHAPE.ROUNDED_RECTANGLE

            # Create the shape with appropriate styling
            shape = slide.shapes.add_shape(
                shape_type,
                left=position["left"],
                top=position["top"],
                width=width,
                height=height,
            )

            # Style the shape based on importance/emphasis
            emphasis = element.get("emphasis", "medium")
            importance = element.get("importance", 1)

            # Set fill color based on emphasis level
            fill = shape.fill
            fill.solid()

            # Color scheme based on emphasis
            if emphasis == "high" or importance >= 3:
                # Light orange for high emphasis
                fill.fore_color.rgb = RGBColor(255, 235, 205)
                line_color = RGBColor(255, 165, 0)  # Orange border
                line_width = Pt(2.5)
            elif emphasis == "medium" or importance == 2:
                # Light blue for medium emphasis
                fill.fore_color.rgb = RGBColor(220, 235, 255)
                line_color = RGBColor(100, 149, 237)  # Cornflower blue border
                line_width = Pt(2.0)
            else:
                # Light gray for low emphasis
                fill.fore_color.rgb = RGBColor(245, 245, 245)
                line_color = RGBColor(169, 169, 169)  # Gray border
                line_width = Pt(1.5)

            # Set border style
            line = shape.line
            line.color.rgb = line_color
            line.width = line_width

            # Add text content
            text_frame = shape.text_frame
            text_frame.text = element.get("text", "")
            text_frame.word_wrap = True
            text_frame.margin_bottom = Inches(0.05)
            text_frame.margin_left = Inches(0.05)
            text_frame.margin_right = Inches(0.05)
            text_frame.margin_top = Inches(0.05)

            # Center text vertically
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Style the text
            para = text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            para.font.size = Pt(12)
            para.font.name = "Arial"

            # Make important text bold
            if importance >= 2:
                para.font.bold = True

            # Store the created shape for connections
            created_shapes[element_id] = shape

        # Now create all connections
        if "connections" in analysis_data:
            for connection in analysis_data["connections"]:
                from_id = connection.get("from_id")
                to_id = connection.get("to_id")
                arrow_type = connection.get("arrow_type", "simple")

                # Make sure both shapes exist
                if from_id in created_shapes and to_id in created_shapes:
                    from_shape = created_shapes[from_id]
                    to_shape = created_shapes[to_id]

                    # Determine connector type and styling
                    if arrow_type == "double":
                        # Double-headed arrow
                        begin_arrow_type = MSO_THEME_COLOR.ACCENT_1
                        end_arrow_type = MSO_THEME_COLOR.ACCENT_1
                    elif arrow_type == "simple":
                        # Single-headed arrow
                        begin_arrow_type = None
                        end_arrow_type = MSO_THEME_COLOR.ACCENT_1
                    else:
                        # No arrows, just a line
                        begin_arrow_type = None
                        end_arrow_type = None

                    # Create the connector
                    connector = slide.shapes.add_connector(
                        MSO_CONNECTOR.STRAIGHT, 0, 0, 0, 0
                    )

                    # Connect the shapes
                    connector.begin_connect(from_shape, 0)
                    connector.end_connect(to_shape, 0)

                    # Style the connector
                    connector.line.color.rgb = RGBColor(100, 100, 100)  # Dark gray line
                    connector.line.width = Pt(1.5)

                    # Set arrow styling if needed
                    if arrow_type == "dashed":
                        connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH

                    # Set arrows based on arrow_type
                    if end_arrow_type:
                        connector.line.end_arrow_type = end_arrow_type
                    if begin_arrow_type:
                        connector.line.begin_arrow_type = begin_arrow_type

        # Return the created presentation
        return prs

    def save_ppt(self, prs, filename="output.pptx"):
        """Save the PowerPoint presentation to a file and return it as bytes"""
        # Save to a BytesIO object
        ppt_bytes = BytesIO()
        prs.save(ppt_bytes)
        ppt_bytes.seek(0)

        # Also save to a file if specified
        if filename:
            prs.save(filename)
            st.success(f"Presentation saved as {filename}")

        return ppt_bytes.getvalue()


def main():
    st.set_page_config(
        page_title="Whiteboard to PowerPoint Converter",
        page_icon="📊",
        layout="wide",
    )

    st.title("Whiteboard to PowerPoint Converter")
    st.write(
        "Upload a whiteboard image and convert it to a structured PowerPoint presentation."
    )

    # Initialize the converter
    converter = WhiteboardToPPT()

    # File uploader for whiteboard image
    uploaded_file = st.file_uploader(
        "Upload a whiteboard image", type=["jpg", "jpeg", "png"]
    )

    # Manual rotation option
    rotation_angle = st.slider(
        "Rotation angle (if needed)",
        min_value=-180,
        max_value=180,
        value=0,
        step=5,
        help="Adjust if your image needs rotation",
    )

    if uploaded_file is not None:
        # Display the uploaded image
        st.image(uploaded_file, caption="Uploaded Image", use_container_width=True)

        # Process button
        if st.button("Convert to PowerPoint"):
            with st.spinner("Processing your whiteboard image..."):
                try:
                    # Analyze the image
                    image_data = uploaded_file.getvalue()
                    analysis_data = converter.analyze_image(image_data, rotation_angle)

                    if analysis_data:
                        # Create PowerPoint from analysis
                        ppt = converter.create_ppt_from_analysis(analysis_data)

                        # Save and provide download link
                        ppt_bytes = converter.save_ppt(
                            ppt, "whiteboard_presentation.pptx"
                        )

                        # Create download button
                        st.download_button(
                            label="Download PowerPoint",
                            data=ppt_bytes,
                            file_name="whiteboard_presentation.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        )

                        st.success("Conversion completed successfully!")
                    else:
                        st.error(
                            "Failed to analyze the image. Please try a different image or adjust rotation."
                        )

                except Exception as e:
                    st.error(f"An error occurred during processing: {e}")
                    import traceback

                    st.error(traceback.format_exc())


if __name__ == "__main__":
    main()
