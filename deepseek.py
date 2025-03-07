import base64
import io
import json
import logging
import os
from io import BytesIO

import requests
import streamlit as st
from dotenv import load_dotenv
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Load environment variables for API keys
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class WhiteboardToPPT:
    def __init__(self):
        # API configuration - use environment variable or Streamlit secrets
        self.vision_api_key = os.getenv("OPENAI_API_KEY")
        self.vision_api_url = "https://api.openai.com/v1/chat/completions"

        # PowerPoint settings
        self.slide_width = Inches(10)
        self.slide_height = Inches(7.5)

        # Position mapping for elements
        self.position_map = {
            "top-left": {"x": 1.2, "y": 1.0},
            "top-center": {"x": 4.0, "y": 1.0},
            "top-right": {"x": 7.0, "y": 1.0},
            "middle-left": {"x": 1.2, "y": 2.8},
            "middle-center": {"x": 4.0, "y": 2.8},
            "middle-right": {"x": 7.0, "y": 2.8},
            "bottom-left": {"x": 1.2, "y": 4.6},
            "bottom-center": {"x": 4.0, "y": 4.6},
            "bottom-right": {"x": 7.0, "y": 4.6},
            "far-right-top": {"x": 8.5, "y": 1.0},
            "far-right-middle": {"x": 8.5, "y": 2.8},
            "far-right-bottom": {"x": 8.5, "y": 4.6},
            "top-left-2": {"x": 2.7, "y": 1.0},
            "middle-left-2": {"x": 2.7, "y": 2.8},
            "bottom-left-2": {"x": 2.7, "y": 4.6},
        }

    def analyze_image(self, image_data):
        """
        Use AI vision model to analyze the whiteboard image and extract information.
        Returns structured data about text content and layout.
        """
        try:
            # Convert image to base64
            image_base64 = base64.b64encode(image_data).decode("utf-8")

            # Prepare prompt for the vision model
            prompt = """
            Analyze this whiteboard image and extract:
            1. All text content with exact wording (avoid duplications).
            2. The spatial layout (what text is in boxes, what's connected to what).
            3. The hierarchical structure (main sections, subsections).
            
            Format your response as JSON with these sections:
            - phases: List of main phase headings at the top.
            - sections: List of vertical sections on the left side.
            - elements: List of all boxes/elements with their text, position, and type.
            - connections: List of connections between elements.
            """

            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {self.vision_api_key}",
            }

            payload = {
                "model": "gpt-4o",
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/jpeg;base64,{image_base64}"
                                },
                            },
                        ],
                    }
                ],
                "max_tokens": 4000,
            }

            logger.info("Analyzing image with AI vision model...")
            response = requests.post(self.vision_api_url, headers=headers, json=payload)

            if response.status_code != 200:
                logger.error(f"API Error: {response.status_code}")
                logger.error(response.text)
                return None

            result = response.json()

            if "choices" not in result or len(result["choices"]) == 0:
                logger.error("Unexpected API response format")
                return None

            content = result["choices"][0]["message"]["content"]
            return self._extract_json_from_text(content)

        except Exception as e:
            logger.error(f"Error analyzing image: {e}")
            return None

    def _extract_json_from_text(self, text):
        """
        Extract JSON content from potentially mixed text.
        """
        start_idx = text.find("{")
        end_idx = text.rfind("}")
        if start_idx >= 0 and end_idx > start_idx:
            return json.loads(text[start_idx : end_idx + 1])
        return {}

    def create_ppt_from_analysis(self, analysis_data):
        """
        Create a PowerPoint slide based on the AI-analyzed content.
        """
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide

        # Add phases, sections, elements, and connections
        self._add_phases(slide, analysis_data.get("phases", []))
        self._add_sections(slide, analysis_data.get("sections", []))
        element_shapes = self._add_elements(slide, analysis_data.get("elements", []))
        self._add_connections(
            slide, element_shapes, analysis_data.get("connections", [])
        )

        # Save the presentation to a BytesIO object
        pptx_io = BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        return pptx_io

    def _add_phases(self, slide, phases):
        """
        Add phase headers at the top of the slide.
        """
        phase_width = Inches(2)
        start_x = Inches(1)
        for i, phase in enumerate(phases):
            phase_x = start_x + (i * phase_width)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left=phase_x,
                top=Inches(0.2),
                width=phase_width,
                height=Inches(0.5),
            )
            shape.line.color.rgb = RGBColor(0, 0, 0)
            shape.fill.background()
            tf = shape.text_frame
            tf.text = phase
            p = tf.paragraphs[0]
            p.font.size = Pt(20)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

    def _add_sections(self, slide, sections):
        """
        Add vertical section labels on the left side of the slide.
        """
        for i, section in enumerate(sections):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left=Inches(0.2),
                top=Inches(1 + i * 1.5),
                width=Inches(0.8),
                height=Inches(1),
            )
            shape.line.color.rgb = RGBColor(0, 0, 0)
            shape.fill.background()
            tf = shape.text_frame
            tf.text = section
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            shape.rotation = 270
            p = tf.paragraphs[0]
            p.font.size = Pt(14)

    def _add_elements(self, slide, elements):
        """
        Add elements (boxes/text) to the slide.
        """
        element_shapes = {}
        for element in elements:
            # Get position configuration or use a default position
            position_key = element.get("position", "middle-center")
            pos_config = self.position_map.get(
                position_key, {"x": 4.0, "y": 3.5}
            )  # Default to middle-center

            # Ensure pos_config is a dictionary with "x" and "y" keys
            if (
                not isinstance(pos_config, dict)
                or "x" not in pos_config
                or "y" not in pos_config
            ):
                pos_config = {"x": 4.0, "y": 3.5}  # Fallback to middle-center

            # Convert coordinates to Inches
            left = Inches(pos_config["x"])
            top = Inches(pos_config["y"])

            # Calculate width and height based on text length
            text_length = len(element.get("text", ""))
            width = min(Inches(3.5), max(Inches(1.5), Inches(text_length / 12)))
            height = Inches(0.8)

            # Create shape for the element
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, height
            )
            shape.line.color.rgb = RGBColor(0, 0, 0)  # Black outline
            shape.fill.background()  # Transparent fill

            # Add text to the shape
            tf = shape.text_frame
            tf.text = element.get("text", "")
            p = tf.paragraphs[0]
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER

            # Store shape information for connections
            element_id = element.get("id", f"elem_{id(element)}")
            element_shapes[element_id] = {
                "shape": shape,
                "left": left,
                "top": top,
                "width": width,
                "height": height,
            }

        return element_shapes

    def _add_connections(self, slide, element_shapes, connections):
        """
        Add connection arrows between elements.
        Skips invalid connections and logs warnings.
        """
        for conn in connections:
            # Skip if the connection is missing required keys
            if (
                not isinstance(conn, dict)
                or "from_id" not in conn
                or "to_id" not in conn
            ):
                logger.warning(f"Skipping invalid connection: {conn}")
                continue

            from_id = conn["from_id"]
            to_id = conn["to_id"]

            # Skip if either shape doesn't exist
            if from_id not in element_shapes or to_id not in element_shapes:
                logger.warning(
                    f"Skipping connection: {from_id} → {to_id} (shape not found)"
                )
                continue

            from_shape = element_shapes[from_id]
            to_shape = element_shapes[to_id]

            # Calculate optimal connection points
            start_x = from_shape["left"] + from_shape["width"]
            start_y = from_shape["top"] + (from_shape["height"] / 2)
            end_x = to_shape["left"]
            end_y = to_shape["top"] + (to_shape["height"] / 2)

            # Create arrow connector
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y
            )
            connector.line.color.rgb = RGBColor(0, 0, 0)  # Black arrow
            connector.line.width = Pt(1.5)
            connector.line.dash_style = MSO_LINE_DASH_STYLE.SOLID

    def process_image_to_ppt(self, image_data):
        """
        Main method to process image and create PowerPoint.
        """
        analysis_data = self.analyze_image(image_data)
        if not analysis_data:
            return None
        return self.create_ppt_from_analysis(analysis_data)


def main():
    st.set_page_config(page_title="Whiteboard to PowerPoint Converter", page_icon="📊")
    st.title("Whiteboard to PowerPoint Converter")

    # Check for API key
    api_key = os.getenv("OPENAI_API_KEY") or st.secrets.get("OPENAI_API_KEY", "")
    if not api_key:
        st.warning("⚠️ OpenAI API key not found. Please add it to your .env file.")
        api_key = st.text_input("Enter your OpenAI API key:", type="password")
        if api_key:
            os.environ["OPENAI_API_KEY"] = api_key

    # File uploader
    uploaded_file = st.file_uploader(
        "Upload whiteboard image", type=["jpg", "jpeg", "png"]
    )
    if uploaded_file is not None:
        image = Image.open(uploaded_file)
        st.image(image, caption="Uploaded Whiteboard Image", use_container_width=True)

        if st.button("Convert to PowerPoint"):
            with st.spinner("Processing..."):
                uploaded_file.seek(0)
                image_data = uploaded_file.read()

                converter = WhiteboardToPPT()
                pptx_data = converter.process_image_to_ppt(image_data)

                if pptx_data:
                    st.success("PowerPoint slide created successfully!")
                    st.download_button(
                        label="Download PowerPoint",
                        data=pptx_data,
                        file_name="whiteboard_to_ppt.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )
                else:
                    st.error(
                        "Failed to create PowerPoint. Please check the logs for details."
                    )


if __name__ == "__main__":
    main()
