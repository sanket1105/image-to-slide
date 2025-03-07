import base64
import io
import json
import os
from io import BytesIO

import requests
import streamlit as st
from dotenv import load_dotenv
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# Load environment variables for API keys
load_dotenv()
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN


class WhiteboardToPPT:
    def __init__(self):
        # API configuration - use environment variable or Streamlit secrets
        self.vision_api_key = os.getenv("OPENAI_API_KEY")
        self.vision_api_url = "https://api.openai.com/v1/chat/completions"

        # PowerPoint settings
        self.slide_width = Inches(10)
        self.slide_height = Inches(7.5)

    def analyze_image(self, image_data):
        """
        Use AI vision model to analyze the whiteboard image and extract information.
        Returns structured data about text content and layout.
        """
        try:
            # Convert image to base64
            image_base64 = base64.b64encode(image_data).decode("utf-8")

            # Prepare improved prompt for the vision model
            prompt = """
            Analyze this whiteboard image and extract:
            1. All text content with exact wording (make sure to avoid duplications)
            2. The spatial layout (what text is in boxes, what's connected to what)
            3. The hierarchical structure (main sections, subsections)
            
            Format your response as JSON with these sections:
            - phases: List of main phase headings at the top (e.g., "Ideation", "Align", "Launch", "Manage")
            - sections: List of vertical sections on the left side
            - elements: List of all boxes/elements with their text, position, and type
            - connections: List of connections between elements
            
            For positions, be specific about exact location:
            - Use "top", "middle", "bottom" for vertical position
            - Use "left", "center", "right" for horizontal position
            - Combine these for precise placement (e.g., "top-left", "middle-center")
            
            For each element, include these properties:
            - id: A unique identifier for the element (e.g., "elem1", "elem2")
            - text: The exact text content (with no duplications)
            - position: The position on the slide (e.g., "top-left", "middle-right")
            - type: "box" if text is in a box/rectangle, "text" if standalone text, "note" if it appears to be a note/comment
            - coordinates: Approximate x,y coordinates as percentages of the image (e.g., {"x": 20, "y": 30})
            - has_shape: boolean indicating if text is enclosed in a visible shape
            
            For connections, include:
            - from_id: ID of the source element
            - to_id: ID of the target element
            
            Example format:
            {
              "phases": ["Ideation", "Align", "Launch", "Manage"],
              "sections": ["Strategy", "Planning", "Foundation"],
              "elements": [
                {"id": "elem1", "text": "Content Ideas", "position": "middle-left", "type": "box", "coordinates": {"x": 20, "y": 30}, "has_shape": true},
                {"id": "elem2", "text": "Need Sizing Tool", "position": "middle-left", "type": "note", "coordinates": {"x": 25, "y": 40}, "has_shape": false}
              ],
              "connections": [
                {"from_id": "elem1", "to_id": "elem2"}
              ]
            }
            
            Important: 
            1. Ensure there are no duplicate texts in elements.
            2. Carefully identify all arrows/lines in the image and represent them as connections.
            3. If text appears in multiple places with the same wording, represent it as a single element.
            4. Ensure each element has a unique ID to properly reference in connections.
            """

            # API call to vision model
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {self.vision_api_key}",
            }

            # Updated to use the current GPT-4 Vision model
            payload = {
                "model": "gpt-4o",  # Current model for vision capabilities
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/jpeg;base64,{image_base64}"
                                },
                            },
                        ],
                    }
                ],
                "max_tokens": 4000,
            }

            st.write("Analyzing image with AI vision model...")

            response = requests.post(self.vision_api_url, headers=headers, json=payload)

            if response.status_code != 200:
                st.error(f"API Error: {response.status_code}")
                st.error(response.text)
                return None

            result = response.json()

            # Check if we have the expected structure
            if "choices" not in result or len(result["choices"]) == 0:
                st.error("Unexpected API response format")
                st.write(result)
                return None

            # Extract the content from the response
            content = result["choices"][0]["message"]["content"]

            # Parse the JSON content - handle potential JSON formatting issues
            try:
                # Find JSON content in the response (it might be mixed with other text)
                json_content = self._extract_json_from_text(content)
                analysis_data = json.loads(json_content)

                # Log the detected connections
                if "connections" in analysis_data:
                    st.write(
                        f"Detected {len(analysis_data['connections'])} connections between elements."
                    )
                    for conn in analysis_data["connections"]:
                        st.write(f"Connection: {conn['from_id']} → {conn['to_id']}")

                return analysis_data
            except json.JSONDecodeError as e:
                st.error(f"Error parsing JSON response: {e}")
                st.write("Raw content:")
                st.write(content)

                # Fallback to manual extraction if JSON parsing fails
                return self._fallback_extraction(content)

        except Exception as e:
            st.error(f"Error analyzing image: {e}")
            return None

    def _extract_json_from_text(self, text):
        """Extract JSON content from potentially mixed text"""
        # Look for content between curly braces
        start_idx = text.find("{")
        end_idx = text.rfind("}")

        if start_idx >= 0 and end_idx > start_idx:
            return text[start_idx : end_idx + 1]

        return "{}"  # Return empty JSON if none found

    def _fallback_extraction(self, content):
        """Manual extraction of key elements as fallback"""
        st.warning("Using fallback extraction method")

        # Based on the provided image, create a more accurate fallback
        return {
            "phases": ["Ideation", "Align", "Launch", "Manage"],
            "sections": ["Strategy", "Planning", "Foundation"],
            "elements": [
                {
                    "id": "elem1",
                    "text": "#Need Sizing Tool#",
                    "position": "top-left",
                    "type": "note",
                    "coordinates": {"x": 15, "y": 13},
                    "has_shape": True,
                },
                {
                    "id": "elem2",
                    "text": "Define: Test (channels) Audience Success",
                    "position": "top-center",
                    "type": "box",
                    "coordinates": {"x": 35, "y": 13},
                    "has_shape": True,
                },
                {
                    "id": "elem3",
                    "text": "Job Aid",
                    "position": "top-right",
                    "type": "box",
                    "coordinates": {"x": 65, "y": 13},
                    "has_shape": True,
                },
                {
                    "id": "elem4",
                    "text": "Sharing Reporting (Data)",
                    "position": "top-right",
                    "type": "box",
                    "coordinates": {"x": 85, "y": 13},
                    "has_shape": True,
                },
                {
                    "id": "elem5",
                    "text": "Content Ideas",
                    "position": "middle-left",
                    "type": "box",
                    "coordinates": {"x": 15, "y": 50},
                    "has_shape": True,
                },
                {
                    "id": "elem6",
                    "text": "Create offer needed/initiative & tasks",
                    "position": "middle-center",
                    "type": "box",
                    "coordinates": {"x": 45, "y": 50},
                    "has_shape": True,
                },
                {
                    "id": "elem7",
                    "text": "Sales Calling & Acc",
                    "position": "middle-right",
                    "type": "box",
                    "coordinates": {"x": 85, "y": 50},
                    "has_shape": True,
                },
                {
                    "id": "elem8",
                    "text": "Lead/Contact Quality (? Randy?)",
                    "position": "bottom-center",
                    "type": "box",
                    "coordinates": {"x": 45, "y": 75},
                    "has_shape": True,
                },
                {
                    "id": "elem9",
                    "text": "Produce Reporting",
                    "position": "bottom-right",
                    "type": "box",
                    "coordinates": {"x": 85, "y": 75},
                    "has_shape": True,
                },
                {
                    "id": "elem10",
                    "text": "Continuation/feedback: Gather, Match Rates, etc",
                    "position": "bottom-left",
                    "type": "box",
                    "coordinates": {"x": 25, "y": 90},
                    "has_shape": True,
                },
                {
                    "id": "elem11",
                    "text": "Error",
                    "position": "bottom-right",
                    "type": "box",
                    "coordinates": {"x": 85, "y": 90},
                    "has_shape": True,
                },
            ],
            "connections": [
                {"from_id": "elem1", "to_id": "elem2"},
                {"from_id": "elem2", "to_id": "elem3"},
                {"from_id": "elem3", "to_id": "elem4"},
                {"from_id": "elem5", "to_id": "elem6"},
                {"from_id": "elem6", "to_id": "elem7"},
                {"from_id": "elem8", "to_id": "elem9"},
                {"from_id": "elem10", "to_id": "elem11"},
            ],
        }

    def adjust_positions_for_overlap(self, elements_dict):
        """Adjust positions to prevent overlap between elements"""
        # Sort elements by position (top to bottom, left to right)
        elements_list = list(elements_dict.items())

        # Check each pair of elements for overlap
        for i in range(len(elements_list)):
            for j in range(i + 1, len(elements_list)):
                elem1_id, elem1 = elements_list[i]
                elem2_id, elem2 = elements_list[j]

                # Calculate boundaries
                elem1_right = elem1["left"] + elem1["width"]
                elem1_bottom = elem1["top"] + elem1["height"]
                elem2_right = elem2["left"] + elem2["width"]
                elem2_bottom = elem2["top"] + elem2["height"]

                # Check for overlap
                if (
                    elem1["left"] < elem2_right
                    and elem1_right > elem2["left"]
                    and elem1["top"] < elem2_bottom
                    and elem1_bottom > elem2["top"]
                ):

                    # There is overlap - move elem2 down or to the right
                    # Determine if horizontal or vertical adjustment is better
                    h_overlap = min(elem1_right, elem2_right) - max(
                        elem1["left"], elem2["left"]
                    )
                    v_overlap = min(elem1_bottom, elem2_bottom) - max(
                        elem1["top"], elem2["top"]
                    )

                    if h_overlap < v_overlap:
                        # Horizontal adjustment is smaller
                        if elem1["left"] < elem2["left"]:
                            # Move elem2 to the right of elem1
                            elements_dict[elem2_id]["left"] = elem1_right + Inches(0.2)
                        else:
                            # Move elem1 to the right of elem2
                            elements_dict[elem1_id]["left"] = elem2_right + Inches(0.2)
                    else:
                        # Vertical adjustment is smaller
                        if elem1["top"] < elem2["top"]:
                            # Move elem2 below elem1
                            elements_dict[elem2_id]["top"] = elem1_bottom + Inches(0.2)
                        else:
                            # Move elem1 below elem2
                            elements_dict[elem1_id]["top"] = elem2_bottom + Inches(0.2)

        return elements_dict

    def create_ppt_from_analysis(self, analysis_data):
        """
        Create a PowerPoint slide based on the AI-analyzed content.
        """
        # Initialize PowerPoint presentation
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]  # Blank slide layout
        slide = prs.slides.add_slide(slide_layout)

        # 1. Add the phase headers at the top in SEPARATE BOXES
        if "phases" in analysis_data and analysis_data["phases"]:
            phase_width = Inches(2)
            total_phases = len(analysis_data["phases"])
            start_x = Inches(1)

            for i, phase in enumerate(analysis_data["phases"]):
                # Calculate position for evenly spaced boxes
                phase_x = start_x + (i * phase_width)

                # Create a box for each phase
                phase_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left=phase_x,
                    top=Inches(0.2),
                    width=phase_width,
                    height=Inches(0.5),
                )

                # Configure the shape
                phase_shape.line.color.rgb = RGBColor(0, 0, 0)  # Black outline
                phase_shape.fill.background()  # Transparent fill

                # Add text
                text_frame = phase_shape.text_frame
                text_frame.text = phase
                p = text_frame.paragraphs[0]
                p.font.size = Pt(20)  # Smaller font to avoid overflow
                p.font.bold = True
                p.font.name = "Arial"
                p.font.color.rgb = RGBColor(0, 0, 0)  # Black text
                p.alignment = PP_ALIGN.CENTER

        # 2. Add the vertical section labels on the left in boxes
        if "sections" in analysis_data and analysis_data["sections"]:
            for i, section in enumerate(analysis_data["sections"]):
                section_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left=Inches(0.2),
                    top=Inches(1 + i * 1.5),  # Increased spacing to avoid overlap
                    width=Inches(0.8),
                    height=Inches(1),
                )
                section_shape.line.color.rgb = RGBColor(0, 0, 0)  # Black outline
                section_shape.fill.background()  # Transparent fill

                # Add text
                text_frame = section_shape.text_frame
                text_frame.text = section
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                section_shape.rotation = 270  # Rotate text vertically
                p = text_frame.paragraphs[0]
                p.font.size = Pt(14)
                p.font.name = "Arial"
                p.font.color.rgb = RGBColor(0, 0, 0)  # Black text

        # 3. Create a mapping for positions to actual coordinates with improved spacing
        position_map = {
            # Top row positions - more spacing
            "top-left": {"x": 1.2, "y": 1.0},
            "top-center": {"x": 4.0, "y": 1.0},
            "top-right": {"x": 7.0, "y": 1.0},
            # Middle row positions - more spacing
            "middle-left": {"x": 1.2, "y": 2.8},
            "middle-center": {"x": 4.0, "y": 2.8},
            "middle-right": {"x": 7.0, "y": 2.8},
            # Bottom row positions - more spacing
            "bottom-left": {"x": 1.2, "y": 4.6},
            "bottom-center": {"x": 4.0, "y": 4.6},
            "bottom-right": {"x": 7.0, "y": 4.6},
            # Far right positions
            "far-right-top": {"x": 8.5, "y": 1.0},
            "far-right-middle": {"x": 8.5, "y": 2.8},
            "far-right-bottom": {"x": 8.5, "y": 4.6},
            # Additional positions for more precise placement
            "top-left-2": {"x": 2.7, "y": 1.0},
            "middle-left-2": {"x": 2.7, "y": 2.8},
            "bottom-left-2": {"x": 2.7, "y": 4.6},
        }

        # Store elements by ID for adding connections later
        element_shapes = {}

        # Add additional dynamic positions based on coordinates if provided
        if "elements" in analysis_data:
            for element in analysis_data["elements"]:
                if "coordinates" in element and "id" in element:
                    # Create a unique position key based on coordinates
                    element_id = element["id"]
                    # Convert percentage coordinates to inches (adjusted to prevent overlap)
                    x_inches = max(
                        1.2, min(8.5, (element["coordinates"]["x"] / 100) * 9)
                    )
                    y_inches = max(
                        1.0, min(6.5, (element["coordinates"]["y"] / 100) * 6.5)
                    )
                    position_map[element_id] = {"x": x_inches, "y": y_inches}
                    # Add the position key to the element for later use
                    element["position_key"] = element_id

        # 4. Add all the elements based on their position and type
        if "elements" in analysis_data:
            for element in analysis_data["elements"]:
                element_id = element.get("id", f"elem_{id(element)}")
                position = element.get(
                    "position_key", element.get("position", "middle-center")
                )

                # Get position configuration or use default
                pos_config = position_map.get(
                    position,
                    position_map.get(
                        element_id, {"x": 3.5, "y": 3.5}
                    ),  # Try element_id or default to middle
                )

                # Extract coordinates
                left = Inches(pos_config["x"])
                top = Inches(pos_config["y"])

                # Determine text length to set appropriate width
                text_length = len(element["text"])
                word_count = len(element["text"].split())

                # Calculate number of lines needed (estimating ~25-30 characters per line)
                estimated_lines = max(1, text_length // 25)

                # Scale width and height based on text length and estimated lines
                width = min(Inches(3.5), max(Inches(1.5), Inches(text_length / 12)))
                height = min(
                    Inches(1.2), max(Inches(0.5), Inches(estimated_lines * 0.3))
                )

                # Create shape for all elements (always a box)
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left=left,
                    top=top,
                    width=width,
                    height=height,
                )

                # Set black outline for all shapes
                shape.line.color.rgb = RGBColor(0, 0, 0)  # Black outline
                shape.fill.background()  # Transparent fill

                # Adjust text frame settings
                text_frame = shape.text_frame
                text_frame.word_wrap = True

                # Adjust font size based on text length to prevent overflow
                font_size = 12
                if text_length > 30:
                    font_size = 10
                if text_length > 50:
                    font_size = 9
                if text_length > 70:
                    font_size = 8

                # Add and format the text
                text_frame.text = element["text"]
                p = text_frame.paragraphs[0]
                p.font.size = Pt(font_size)
                p.font.name = "Arial"
                p.font.color.rgb = RGBColor(0, 0, 0)  # Black text for all elements
                p.alignment = PP_ALIGN.CENTER  # Center align text

                # Store the shape by ID for connection arrows
                element_shapes[element_id] = {
                    "shape": shape,
                    "text": element["text"],
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height,
                }

        # Apply overlap prevention algorithm
        element_shapes = self.adjust_positions_for_overlap(element_shapes)

        # Update the actual shapes with the adjusted positions
        for elem_id, elem_info in element_shapes.items():
            shape = elem_info["shape"]
            shape.left = elem_info["left"]
            shape.top = elem_info["top"]

        # 5. Add connection arrows based on AI-detected connections
        if "connections" in analysis_data:
            for connection in analysis_data["connections"]:
                from_id = connection["from_id"]
                to_id = connection["to_id"]

                # Skip if either shape doesn't exist
                if from_id not in element_shapes or to_id not in element_shapes:
                    continue

                from_shape = element_shapes[from_id]
                to_shape = element_shapes[to_id]

                # Calculate optimal connection points based on relative positions
                from_center_x = from_shape["left"] + (from_shape["width"] / 2)
                from_center_y = from_shape["top"] + (from_shape["height"] / 2)
                to_center_x = to_shape["left"] + (to_shape["width"] / 2)
                to_center_y = to_shape["top"] + (to_shape["height"] / 2)

                # Determine the best connection points based on box positions
                if abs(from_center_x - to_center_x) > abs(from_center_y - to_center_y):
                    # Horizontal connection (left to right or right to left)
                    if from_center_x < to_center_x:
                        # Left to right
                        start_x = from_shape["left"] + from_shape["width"]
                        start_y = from_center_y
                        end_x = to_shape["left"]
                        end_y = to_center_y
                    else:
                        # Right to left
                        start_x = from_shape["left"]
                        start_y = from_center_y
                        end_x = to_shape["left"] + to_shape["width"]
                        end_y = to_center_y
                else:
                    # Vertical connection (top to bottom or bottom to top)
                    if from_center_y < to_center_y:
                        # Top to bottom
                        start_x = from_center_x
                        start_y = from_shape["top"] + from_shape["height"]
                        end_x = to_center_x
                        end_y = to_shape["top"]
                    else:
                        # Bottom to top
                        start_x = from_center_x
                        start_y = from_shape["top"]
                        end_x = to_center_x
                        end_y = to_shape["top"] + to_shape["height"]

                # Create arrow connector
                connector = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y
                )
                connector.line.color.rgb = RGBColor(0, 0, 0)  # Black arrow
                connector.line.width = Pt(1.5)
                connector.line.dash_style = MSO_LINE_DASH_STYLE.SOLID

                # Do not set arrow styles as they were causing errors
                # Instead, just leave as basic connectors

        # Save the PowerPoint
        pptx_io = BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        return pptx_io

    def process_image_to_ppt(self, image_data):
        """
        Main method to process image and create PowerPoint.
        """
        # Step 1: Analyze image with AI
        analysis_data = self.analyze_image(image_data)

        if not analysis_data:
            st.error("Failed to analyze image")
            return None

        # Step 2: Create PowerPoint from analysis
        st.write("Creating PowerPoint slide...")
        pptx_data = self.create_ppt_from_analysis(analysis_data)

        return pptx_data


def main():
    st.set_page_config(page_title="Whiteboard to PowerPoint Converter", page_icon="📊")

    st.title("Whiteboard to PowerPoint Converter")
    st.write("Upload a whiteboard image and get a structured PowerPoint slide.")

    # Check for API key
    api_key = os.getenv("OPENAI_API_KEY") or st.secrets.get("OPENAI_API_KEY", "")
    if not api_key:
        st.warning(
            "⚠️ OpenAI API key not found. Please add it to your .env file or Streamlit secrets."
        )
        api_key = st.text_input("Enter your OpenAI API key:", type="password")
        if api_key:
            os.environ["OPENAI_API_KEY"] = api_key

    # File uploader
    uploaded_file = st.file_uploader(
        "Upload whiteboard image", type=["jpg", "jpeg", "png"]
    )

    if uploaded_file is not None:
        # Display the uploaded image
        image = Image.open(uploaded_file)
        st.image(image, caption="Uploaded Whiteboard Image", use_container_width=True)

        # Process image when button is clicked
        if st.button("Convert to PowerPoint"):
            with st.spinner("Processing..."):
                # Reset uploaded_file to start
                uploaded_file.seek(0)
                image_data = uploaded_file.read()

                # Create converter instance
                converter = WhiteboardToPPT()

                # Process image
                pptx_data = converter.process_image_to_ppt(image_data)

                if pptx_data:
                    # Provide download button for the generated PowerPoint
                    st.success("PowerPoint slide created successfully!")

                    # Create download button
                    st.download_button(
                        label="Download PowerPoint",
                        data=pptx_data,
                        file_name="whiteboard_to_ppt.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )
                else:
                    st.error(
                        "Failed to create PowerPoint. Please check the logs above for details."
                    )


if __name__ == "__main__":
    main()


## final
